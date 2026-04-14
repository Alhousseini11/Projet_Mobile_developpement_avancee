from __future__ import annotations

import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "presentation_assets" / "sprint_6"
OUTPUT_PATH = ROOT / "presentation_sprint_6_demo.pptx"

WIDE_W = 13.333
WIDE_H = 7.5

PRIMARY = RGBColor(220, 38, 38)
PRIMARY_SOFT = RGBColor(254, 226, 226)
DARK = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)

PIL_PRIMARY = (220, 38, 38)
PIL_PRIMARY_DARK = (153, 27, 27)
PIL_DARK = (15, 23, 42)
PIL_PANEL = (30, 41, 59)
PIL_SHELL = (241, 245, 249)
PIL_SOFT = (248, 250, 252)
PIL_WHITE = (255, 255, 255)
PIL_MUTED = (100, 116, 139)
PIL_TEXT = (15, 23, 42)
PIL_SUCCESS = (22, 101, 52)
PIL_SUCCESS_BG = (220, 252, 231)

LOGIN_SCREEN = ROOT / "frontend" / "login-screen.png"
TUTORIALS_SCREEN = ROOT / "frontend" / "tutorials-final.png"

logo_candidates = list((ROOT / "frontend" / "app" / "assets").glob("Logo-centr*.png"))
LOGO_ASSET = logo_candidates[0] if logo_candidates else ROOT / "frontend" / "app" / "assets" / "Logo.png"

STORIES = [
    {
        "number": "01",
        "title": "Logo de l'application",
        "visual_key": "logo_story",
        "track": "mobile",
        "definition": (
            "Le logo a ete centralise dans un composant reutilisable pour assurer une identite visuelle "
            "constante sur les ecrans d'authentification et dans l'univers mobile Garage Mechanic."
        ),
        "points": [
            "Composant unique pour l'image de marque sur login et ecrans hero.",
            "Deux tailles prises en charge pour s'adapter au contexte visuel.",
            "Renforce la coherence graphique sans dupliquer la logique UI.",
        ],
        "subtasks": [
            "creation du composant BrandLogo",
            "integration du logo sur les ecrans d'authentification",
            "gestion de deux formats hero et compact",
        ],
        "code": {
            "path": "frontend/app/components/BrandLogo.vue",
            "anchor": "<GridLayout class=\"brand-logo\" :class=\"size\">",
            "before": 0,
            "after": 18,
            "max_lines": 18,
            "label": "frontend/app/components/BrandLogo.vue",
            "explanation": (
                "Ce composant encapsule le logo et expose une prop size. "
                "L'application reutilise donc le meme visuel sans re-coder le style sur chaque page."
            ),
        },
    },
    {
        "number": "02",
        "title": "Lecture video sur la page tutorielle",
        "visual_key": "tutorial_player_story",
        "track": "mobile",
        "definition": (
            "La consultation d'un tutoriel ne se limite plus a une fiche descriptive : l'utilisateur ouvre "
            "un ecran de lecture integre, reste dans l'application et suit la video directement depuis le parcours tutoriels."
        ),
        "points": [
            "Navigation detail -> lecteur video sans sortir vers le navigateur.",
            "Lecture integree dans un WebView embarque dedie a la video.",
            "Retour utilisateur propre vers l'ecran precedent du parcours tutoriels.",
        ],
        "subtasks": [
            "creation d'un ecran detail tutoriel",
            "creation du lecteur video integre",
            "navigation detail -> lecteur -> retour utilisateur",
        ],
        "code": {
            "path": "frontend/app/components/TutorialVideoPlayer.vue",
            "anchor": "function scheduleQualifiedViewCount() {",
            "before": 0,
            "after": 26,
            "max_lines": 22,
            "label": "frontend/app/components/TutorialVideoPlayer.vue",
            "explanation": (
                "Le lecteur attend quelques secondes de lecture reelle avant d'envoyer la vue au backend. "
                "On evite ainsi de compter un simple clic sur le bouton Lire la video."
            ),
        },
    },
    {
        "number": "03",
        "title": "Console admin pour upload les videos",
        "visual_key": "admin_upload_story",
        "track": "admin",
        "definition": (
            "Un front web admin separe de l'application mobile permet de publier les tutoriels video "
            "depuis un poste de travail, avec televersement de fichier et metadata completes."
        ),
        "points": [
            "Interface web dediee au role ADMIN.",
            "Upload direct d'un fichier video depuis l'ordinateur.",
            "Categorie, difficulte, duree, outils et etapes enregistres avec la video.",
        ],
        "subtasks": [
            "creation du formulaire web d'ajout de tutoriel",
            "construction d'un FormData multipart avec la video",
            "publication du tutoriel directement dans le catalogue client",
        ],
        "code": {
            "path": "admin-web/src/app.ts",
            "anchor": "async function handleTutorialSubmit(event: SubmitEvent) {",
            "before": 0,
            "after": 34,
            "max_lines": 24,
            "label": "admin-web/src/app.ts",
            "explanation": (
                "Le front admin construit un FormData multipart avec le fichier video et toutes les donnees "
                "du tutoriel, puis l'envoie a l'API admin dediee."
            ),
        },
    },
    {
        "number": "04",
        "title": "Routes admin protegees",
        "visual_key": "admin_routes_story",
        "track": "backend",
        "definition": (
            "Les fonctionnalites d'administration reposent sur un namespace backend propre. "
            "Seuls les utilisateurs authentifies avec le role ADMIN peuvent lire les indicateurs, gerer les services et publier les tutoriels."
        ),
        "points": [
            "Toutes les routes admin sont centralisees sous /api/admin.",
            "Protection combinee par authGuard et roleGuard.",
            "Surface API claire pour le dashboard, les services et les tutoriels.",
        ],
        "subtasks": [
            "creation du routeur /api/admin",
            "protection des routes par authentification",
            "restriction supplementaire par role ADMIN",
        ],
        "code": {
            "path": "backend/src/modules/admin/admin.routes.ts",
            "anchor": "router.use(authGuard, roleGuard([Role.ADMIN]));",
            "before": 4,
            "after": 14,
            "max_lines": 18,
            "label": "backend/src/modules/admin/admin.routes.ts",
            "explanation": (
                "Le middleware verrouille tout le routeur avant meme d'entrer dans les handlers. "
                "Un compte client classique ne peut donc pas utiliser les endpoints d'administration."
            ),
        },
    },
    {
        "number": "05",
        "title": "Modification des services cote admin",
        "visual_key": "admin_services_story",
        "track": "admin",
        "definition": (
            "Le catalogue de services est devenu pilotable depuis la console web : l'administrateur peut "
            "creer un service, le modifier, ajuster ses horaires et l'archiver lorsqu'il ne doit plus apparaitre au public."
        ),
        "points": [
            "Edition du libelle, du tarif, de la duree et des creneaux.",
            "Selection d'une carte service pour ouvrir sa fiche detaillee.",
            "Archivage pour retirer un service du catalogue public sans perdre l'historique.",
        ],
        "subtasks": [
            "affichage cliquable des cartes services",
            "edition du service dans le panneau detail",
            "archivage pour retirer le service du catalogue public",
        ],
        "code": {
            "path": "admin-web/src/app.ts",
            "anchor": "async function handleServiceSubmit(event: SubmitEvent) {",
            "before": 0,
            "after": 33,
            "max_lines": 24,
            "label": "admin-web/src/app.ts",
            "explanation": (
                "Le formulaire admin choisit automatiquement entre creation et mise a jour selon le mode courant. "
                "Le meme espace de travail sert donc a enrichir et a maintenir le catalogue."
            ),
        },
    },
    {
        "number": "06",
        "title": "Reinitialisation du mot de passe par courriel",
        "visual_key": "password_reset_story",
        "track": "backend",
        "definition": (
            "La reinitialisation ne repose plus sur un jeton manuel : l'utilisateur demande un code par email, "
            "le recoit via SendGrid et finalise la redefinition du mot de passe depuis l'ecran mobile dedie."
        ),
        "points": [
            "Demande de code depuis l'application mobile.",
            "Envoi du courriel transactionnel par SendGrid.",
            "Saisie du code a 6 chiffres puis definition du nouveau mot de passe.",
        ],
        "subtasks": [
            "demande du code depuis l'ecran mobile oublie mot de passe",
            "generation et envoi du courriel transactionnel",
            "validation du code puis changement du mot de passe",
        ],
        "code": {
            "path": "backend/src/modules/auth/passwordResetMailer.ts",
            "anchor": "function buildPasswordResetEmailContent(payload: {",
            "before": 0,
            "after": 32,
            "max_lines": 24,
            "label": "backend/src/modules/auth/passwordResetMailer.ts",
            "explanation": (
                "Le backend compose un vrai email transactionnel avec le code de reinitialisation et sa date d'expiration. "
                "Le message reste lisible pour l'utilisateur et reutilisable en production."
            ),
        },
    },
    {
        "number": "07",
        "title": "Refonte de la page tutorielle frontend",
        "visual_key": "tutorial_front_story",
        "track": "mobile",
        "definition": (
            "La page tutoriels a ete restructuree pour une navigation plus professionnelle : hero informatif, "
            "moteur de recherche, filtres par categorie et cartes lisibles qui ouvrent une fiche detaillee."
        ),
        "points": [
            "Hero avec statistiques et positionnement plus clair du module video.",
            "Recherche texte et filtres categorie pour reduire le catalogue.",
            "Cartes plus lisibles, detail mieux hierarchise et CTA explicites.",
        ],
        "subtasks": [
            "refonte du hero et des indicateurs visuels",
            "mise en place de la recherche et des filtres categorie",
            "restructuration des cartes et du detail tutoriel",
        ],
        "code": {
            "path": "frontend/app/components/Tutorials.vue",
            "anchor": "function applyFilters() {",
            "before": 0,
            "after": 26,
            "max_lines": 22,
            "label": "frontend/app/components/Tutorials.vue",
            "explanation": (
                "Cette logique applique a la fois le filtre categorie et la recherche texte sur le catalogue. "
                "Le rendu reste reactif et la page conserve une experience de consultation claire."
            ),
        },
    },
]


def ensure_asset_dir() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def load_font(size: int, *, mono: bool = False, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    fonts_dir = Path("C:/Windows/Fonts")
    if mono:
        candidates = [
            fonts_dir / ("consolab.ttf" if bold else "consola.ttf"),
            fonts_dir / ("courbd.ttf" if bold else "cour.ttf"),
        ]
    else:
        candidates = [
            fonts_dir / ("segoeuib.ttf" if bold else "segoeui.ttf"),
            fonts_dir / ("arialbd.ttf" if bold else "arial.ttf"),
            fonts_dir / ("calibrib.ttf" if bold else "calibri.ttf"),
        ]

    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)

    return ImageFont.load_default()


def read_file_lines(relative_path: str) -> list[str]:
    return (ROOT / relative_path).read_text(encoding="utf-8").splitlines()


def trim_common_indent(lines: list[str]) -> list[str]:
    non_empty = [line for line in lines if line.strip()]
    if not non_empty:
        return lines

    indent = min(len(line) - len(line.lstrip(" ")) for line in non_empty)
    return [line[indent:] if len(line) >= indent else line for line in lines]


def extract_snippet(relative_path: str, anchor: str, before: int, after: int, max_lines: int) -> str:
    lines = read_file_lines(relative_path)
    anchor_index = next((index for index, line in enumerate(lines) if anchor in line), None)
    if anchor_index is None:
        raise ValueError(f"Anchor not found in {relative_path}: {anchor}")

    start = max(anchor_index - before, 0)
    end = min(anchor_index + after, len(lines) - 1)
    selected = trim_common_indent(lines[start : end + 1])[:max_lines]
    numbered = [f"{start + offset + 1:>4} | {line.rstrip()}" for offset, line in enumerate(selected)]
    return "\n".join(numbered)


def wrap_code_lines(snippet: str, max_chars: int) -> str:
    wrapped_lines: list[str] = []
    for raw_line in snippet.splitlines():
        if len(raw_line) <= max_chars:
            wrapped_lines.append(raw_line)
            continue

        prefix = raw_line[:7] if "|" in raw_line[:10] else ""
        content = raw_line[7:] if prefix else raw_line
        chunks = textwrap.wrap(
            content,
            width=max_chars - len(prefix),
            replace_whitespace=False,
            drop_whitespace=False,
            break_long_words=False,
            break_on_hyphens=False,
        )
        if not chunks:
            wrapped_lines.append(raw_line)
            continue

        wrapped_lines.append(f"{prefix}{chunks[0]}")
        continuation_prefix = " " * len(prefix)
        for chunk in chunks[1:]:
            wrapped_lines.append(f"{continuation_prefix}{chunk}")
    return "\n".join(wrapped_lines)


def render_code_panel(label: str, snippet: str, destination: Path) -> Path:
    width, height = 1600, 620
    image = Image.new("RGB", (width, height), PIL_DARK)
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((0, 0, width - 1, height - 1), radius=30, fill=PIL_DARK)
    draw.rounded_rectangle((0, 0, width - 1, 76), radius=30, fill=PIL_PANEL)
    draw.text((34, 20), label, font=load_font(30, bold=True), fill=PIL_WHITE)
    draw.text((width - 200, 24), "code", font=load_font(20), fill=(148, 163, 184))

    wrapped = wrap_code_lines(snippet, 94)
    top = 104
    line_height = 28
    for index, line in enumerate(wrapped.splitlines()):
        draw.text((34, top + index * line_height), line, font=load_font(21, mono=True), fill=(226, 232, 240))

    image.save(destination)
    return destination


def fit_image(path: Path, size: tuple[int, int]) -> Image.Image:
    with Image.open(path) as source:
        return ImageOps.fit(source.convert("RGB"), size, method=Image.Resampling.LANCZOS)


def draw_shadow(draw: ImageDraw.ImageDraw, bounds: tuple[int, int, int, int], radius: int = 28) -> None:
    x0, y0, x1, y1 = bounds
    draw.rounded_rectangle((x0 + 12, y0 + 16, x1 + 12, y1 + 16), radius=radius, fill=(221, 228, 235))


def paste_card(base: Image.Image, content: Image.Image, bounds: tuple[int, int, int, int], radius: int = 28) -> None:
    draw = ImageDraw.Draw(base)
    draw_shadow(draw, bounds, radius=radius)
    x0, y0, x1, y1 = bounds
    draw.rounded_rectangle(bounds, radius=radius, fill=PIL_WHITE)
    inner = Image.new("RGB", (x1 - x0 - 28, y1 - y0 - 28), PIL_WHITE)
    inner.paste(content, (0, 0))
    mask = Image.new("L", inner.size, 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, inner.size[0], inner.size[1]), radius=22, fill=255)
    base.paste(inner, (x0 + 14, y0 + 14), mask)


def rounded_panel(
    draw: ImageDraw.ImageDraw,
    bounds: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    radius: int = 26,
    outline: tuple[int, int, int] | None = None,
    width: int = 2,
) -> None:
    if outline:
        draw.rounded_rectangle(bounds, radius=radius, fill=fill, outline=outline, width=width)
    else:
        draw.rounded_rectangle(bounds, radius=radius, fill=fill)


def render_web_shell(draw: ImageDraw.ImageDraw, width: int, height: int, title: str, subtitle: str) -> None:
    rounded_panel(draw, (28, 28, width - 28, height - 28), PIL_WHITE, radius=34)
    rounded_panel(draw, (28, 28, width - 28, 150), PIL_DARK, radius=34)
    rounded_panel(draw, (width - 280, 62, width - 74, 128), (127, 29, 29), radius=22)
    draw.text((64, 58), "GARAGE MECHANIC", font=load_font(22, bold=True), fill=(248, 250, 252))
    draw.text((64, 98), title, font=load_font(38, bold=True), fill=PIL_WHITE)
    draw.text((64, 190), subtitle, font=load_font(22), fill=PIL_MUTED)


def render_logo_story_visual() -> Path:
    destination = ASSET_DIR / "story_01_logo_visual.png"
    width, height = 1200, 1480
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    login_preview = fit_image(LOGIN_SCREEN, (500, 1120))
    logo_preview = fit_image(LOGO_ASSET, (420, 420))

    paste_card(canvas, login_preview, (54, 80, 560, 1240), radius=34)

    rounded_panel(draw, (640, 112, 1120, 420), PIL_DARK, radius=32)
    draw.text((680, 152), "Story branding", font=load_font(22, bold=True), fill=(248, 113, 113))
    draw.text((680, 194), "Identite\nGarage Mechanic", font=load_font(42, bold=True), fill=PIL_WHITE)
    canvas.paste(logo_preview, (680, 480))

    rounded_panel(draw, (640, 940, 1120, 1240), PIL_WHITE, radius=32, outline=(226, 232, 240))
    draw.text((678, 976), "Ce que voit l'utilisateur", font=load_font(28, bold=True), fill=PIL_TEXT)
    items = [
        "logo centre dans le hero",
        "presentation plus professionnelle a la connexion",
        "meme composant reutilise dans l'application",
    ]
    y = 1026
    for item in items:
        rounded_panel(draw, (678, y + 8, 696, y + 26), PIL_PRIMARY, radius=9)
        draw.text((714, y), item, font=load_font(21), fill=PIL_MUTED)
        y += 58

    draw.text((74, 1320), "Capture mobile + zoom du composant de marque", font=load_font(20), fill=PIL_MUTED)
    canvas.save(destination)
    return destination


def render_tutorial_playback_visual() -> Path:
    destination = ASSET_DIR / "story_02_tutorial_playback_visual.png"
    width, height = 1200, 1480
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    tutorials_preview = fit_image(TUTORIALS_SCREEN, (470, 1140))
    paste_card(canvas, tutorials_preview, (48, 72, 530, 1240), radius=34)

    rounded_panel(draw, (596, 70, 1140, 550), PIL_WHITE, radius=34)
    draw.text((632, 108), "Detail du tutoriel", font=load_font(30, bold=True), fill=PIL_TEXT)
    rounded_panel(draw, (632, 160, 1102, 292), PIL_SOFT, radius=24)
    draw.text((660, 184), "vidange", font=load_font(34, bold=True), fill=PIL_TEXT)
    draw.text((660, 228), "Entretien - Niveau moyen - 45 min", font=load_font(21), fill=PIL_MUTED)
    rounded_panel(draw, (632, 330, 1102, 432), PIL_SOFT, radius=22)
    draw.text((658, 360), "Outils requis : cle 12, 15", font=load_font(22), fill=PIL_TEXT)
    rounded_panel(draw, (632, 460, 1102, 518), PIL_PRIMARY, radius=18)
    draw.text((714, 476), "Lire la video dans l'application", font=load_font(22, bold=True), fill=PIL_WHITE)

    rounded_panel(draw, (596, 598, 1140, 1240), PIL_DARK, radius=34)
    draw.text((632, 634), "Lecteur integre", font=load_font(30, bold=True), fill=PIL_WHITE)
    rounded_panel(draw, (632, 690, 1102, 1044), (2, 6, 23), radius=18)
    rounded_panel(draw, (632, 1058, 1102, 1180), PIL_PANEL, radius=18)
    draw.text((654, 1088), "Video embarquee, lecture dans le parcours tutoriels", font=load_font(21), fill=(226, 232, 240))
    rounded_panel(draw, (654, 760, 1080, 952), (17, 24, 39), radius=18)
    rounded_panel(draw, (736, 806, 994, 900), (229, 231, 235), radius=20)
    draw.polygon([(854, 826), (854, 880), (904, 853)], fill=(31, 41, 55))
    draw.text((652, 1120), "La video est lue dans l'application puis comptee apres lecture qualifiee.", font=load_font(20), fill=(148, 163, 184))

    draw.text((78, 1324), "Catalogue mobile reel + detail/lecteur integres", font=load_font(20), fill=PIL_MUTED)
    canvas.save(destination)
    return destination


def render_admin_upload_visual() -> Path:
    destination = ASSET_DIR / "story_03_admin_upload_visual.png"
    width, height = 1480, 1080
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    render_web_shell(
        draw,
        width,
        height,
        "Console admin | Tutoriels video",
        "Upload depuis le poste admin pour alimenter la bibliotheque video client.",
    )

    rounded_panel(draw, (60, 252, 896, 990), PIL_WHITE, radius=30)
    draw.text((96, 292), "Ajouter un tutoriel video", font=load_font(34, bold=True), fill=PIL_TEXT)
    labels = [
        ("Titre", "Vidange atelier - sequence complete"),
        ("Categorie", "entretien"),
        ("Difficulte", "moyen"),
        ("Duree", "45"),
        ("Miniature", "https://placehold.co/..."),
    ]
    y = 350
    for label, value in labels:
        draw.text((96, y), label, font=load_font(20, bold=True), fill=PIL_MUTED)
        rounded_panel(draw, (96, y + 34, 860, y + 102), PIL_SOFT, radius=18)
        draw.text((120, y + 54), value, font=load_font(22), fill=PIL_TEXT)
        y += 118

    draw.text((96, 944), "Etapes, outils et fichier video prepares dans le meme formulaire.", font=load_font(20), fill=PIL_MUTED)

    rounded_panel(draw, (944, 252, 1418, 612), PIL_DARK, radius=30)
    draw.text((980, 290), "Fichier video", font=load_font(30, bold=True), fill=PIL_WHITE)
    rounded_panel(draw, (980, 344, 1382, 548), (17, 24, 39), radius=22, outline=(148, 163, 184), width=3)
    draw.text((1022, 414), "Glisser ou choisir", font=load_font(28, bold=True), fill=PIL_WHITE)
    draw.text((1022, 456), "videoFile.mp4  |  84 Mo", font=load_font(22), fill=(203, 213, 225))
    rounded_panel(draw, (1022, 500, 1248, 546), PIL_SUCCESS_BG, radius=14)
    draw.text((1044, 512), "Format accepte : MP4", font=load_font(18, bold=True), fill=PIL_SUCCESS)

    rounded_panel(draw, (944, 650, 1418, 990), PIL_WHITE, radius=30)
    draw.text((980, 688), "Resultat attendu", font=load_font(30, bold=True), fill=PIL_TEXT)
    cards = [
        "video stockee sur le backend",
        "url publique generee automatiquement",
        "tutorial visible dans le front mobile",
    ]
    y = 754
    for card in cards:
        rounded_panel(draw, (980, y, 1382, y + 70), PIL_SOFT, radius=18)
        rounded_panel(draw, (1002, y + 24, 1018, y + 40), PIL_PRIMARY, radius=8)
        draw.text((1040, y + 18), card, font=load_font(21), fill=PIL_TEXT)
        y += 92

    canvas.save(destination)
    return destination


def render_admin_routes_visual() -> Path:
    destination = ASSET_DIR / "story_04_admin_routes_visual.png"
    width, height = 1480, 1080
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    render_web_shell(
        draw,
        width,
        height,
        "Console admin | Routes protegees",
        "Surface API centralisee pour l'administration et securisee par role.",
    )

    rounded_panel(draw, (60, 250, 450, 990), PIL_DARK, radius=30)
    draw.text((96, 290), "Namespace", font=load_font(24, bold=True), fill=(248, 113, 113))
    draw.text((96, 334), "/api/admin", font=load_font(42, bold=True), fill=PIL_WHITE)
    draw.text((96, 410), "Protection appliquee", font=load_font(22, bold=True), fill=(226, 232, 240))
    draw.text((96, 454), "authGuard + roleGuard([ADMIN])", font=load_font(24), fill=(203, 213, 225))
    draw.text((96, 554), "Ce que cela garantit", font=load_font(22, bold=True), fill=(226, 232, 240))
    benefits = [
        "acces reserve au compte admin",
        "API dediee au back-office",
        "surface de maintenance plus claire",
    ]
    y = 606
    for item in benefits:
        rounded_panel(draw, (96, y + 10, 112, y + 26), PIL_PRIMARY, radius=8)
        draw.text((132, y), item, font=load_font(20), fill=(203, 213, 225))
        y += 62

    rounded_panel(draw, (490, 250, 1418, 990), PIL_WHITE, radius=30)
    draw.text((526, 290), "Routes metier exposees", font=load_font(34, bold=True), fill=PIL_TEXT)
    endpoints = [
        ("GET", "/summary", "charger la synthese du dashboard"),
        ("GET", "/users", "lister les comptes clients"),
        ("GET", "/reservations", "suivre les rendez-vous"),
        ("GET", "/services", "charger le catalogue admin"),
        ("POST", "/services", "creer un nouveau service"),
        ("PUT", "/services/:serviceId", "modifier un service"),
        ("DELETE", "/services/:serviceId", "archiver un service"),
        ("GET", "/tutorials", "charger les tutoriels admin"),
        ("POST", "/tutorials", "televerser une video et publier le tutoriel"),
    ]
    y = 352
    method_colors = {
        "GET": (226, 232, 240),
        "POST": (220, 252, 231),
        "PUT": (254, 240, 138),
        "DELETE": (254, 226, 226),
    }
    for method, route, desc in endpoints:
        rounded_panel(draw, (526, y, 1382, y + 62), PIL_SOFT, radius=18)
        rounded_panel(draw, (550, y + 14, 642, y + 48), method_colors[method], radius=12)
        draw.text((575, y + 18), method, font=load_font(18, bold=True), fill=PIL_TEXT)
        draw.text((674, y + 14), route, font=load_font(22, bold=True), fill=PIL_TEXT)
        draw.text((980, y + 16), desc, font=load_font(18), fill=PIL_MUTED)
        y += 74

    canvas.save(destination)
    return destination


def render_admin_services_visual() -> Path:
    destination = ASSET_DIR / "story_05_admin_services_visual.png"
    width, height = 1480, 1080
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    render_web_shell(
        draw,
        width,
        height,
        "Console admin | Services",
        "Gestion claire du catalogue public, edition detaillee et archivage.",
    )

    rounded_panel(draw, (60, 250, 560, 990), PIL_WHITE, radius=30)
    draw.text((96, 290), "Cartes services", font=load_font(32, bold=True), fill=PIL_TEXT)
    services = [
        ("vidange-premium", "Vidange premium", "79 CAD | 45 min", True),
        ("freins-complets", "Freins complets", "149 CAD | 60 min", False),
        ("diagnostic", "Diagnostic electronique", "59 CAD | 30 min", False),
    ]
    y = 350
    for slug, label, meta, active in services:
        fill = PIL_PRIMARY if active else PIL_SOFT
        text_color = PIL_WHITE if active else PIL_TEXT
        muted_color = (254, 226, 226) if active else PIL_MUTED
        rounded_panel(draw, (96, y, 524, y + 150), fill, radius=24)
        draw.text((126, y + 24), slug, font=load_font(18, bold=True), fill=muted_color)
        draw.text((126, y + 60), label, font=load_font(28, bold=True), fill=text_color)
        draw.text((126, y + 104), meta, font=load_font(21), fill=muted_color if active else PIL_MUTED)
        y += 176

    rounded_panel(draw, (602, 250, 1418, 990), PIL_WHITE, radius=30)
    draw.text((638, 290), "Edition du service selectionne", font=load_font(34, bold=True), fill=PIL_TEXT)
    fields = [
        ("Libelle", "Vidange premium"),
        ("Slug", "vidange-premium"),
        ("Description", "Intervention atelier avec verifications complementaires."),
        ("Duree", "45"),
        ("Prix", "79"),
        ("Horaires", "09:00, 10:30, 13:30"),
    ]
    y = 350
    for label, value in fields:
        draw.text((638, y), label, font=load_font(20, bold=True), fill=PIL_MUTED)
        rounded_panel(draw, (638, y + 34, 1378, y + 102), PIL_SOFT, radius=18)
        draw.text((662, y + 54), value, font=load_font(22), fill=PIL_TEXT)
        y += 118

    rounded_panel(draw, (638, 904, 968, 966), PIL_DARK, radius=18)
    draw.text((718, 922), "Enregistrer", font=load_font(22, bold=True), fill=PIL_WHITE)
    rounded_panel(draw, (1002, 904, 1378, 966), PRIMARY_SOFT_TO_PIL(), radius=18)
    draw.text((1116, 922), "Archiver", font=load_font(22, bold=True), fill=PIL_PRIMARY_DARK)

    canvas.save(destination)
    return destination


def render_password_reset_visual() -> Path:
    destination = ASSET_DIR / "story_06_password_reset_visual.png"
    width, height = 1200, 1480
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    rounded_panel(draw, (48, 70, 556, 1260), PIL_WHITE, radius=34)
    rounded_panel(draw, (48, 70, 556, 306), PIL_DARK, radius=34)
    draw.text((84, 116), "RECUPERATION", font=load_font(20, bold=True), fill=(248, 113, 113))
    draw.text((84, 156), "Mot de passe\noublie", font=load_font(44, bold=True), fill=PIL_WHITE)
    draw.text((84, 264), "Demander un code par email puis\nle saisir dans l'application.", font=load_font(20), fill=(203, 213, 225))
    draw.text((84, 360), "Etape 1", font=load_font(18, bold=True), fill=PIL_PRIMARY)
    draw.text((84, 398), "Recevoir le code", font=load_font(32, bold=True), fill=PIL_TEXT)
    rounded_panel(draw, (84, 470, 520, 548), (236, 253, 245), radius=18)
    draw.text((108, 496), "Lien envoye. Consultez votre boite mail.", font=load_font(18), fill=PIL_SUCCESS)
    draw.text((84, 592), "Adresse email", font=load_font(20, bold=True), fill=PIL_MUTED)
    rounded_panel(draw, (84, 626, 520, 700), PIL_SOFT, radius=18)
    draw.text((108, 650), "alhousseini.adamou@etud.iga.ac.ma", font=load_font(20), fill=PIL_TEXT)
    rounded_panel(draw, (84, 740, 520, 820), PIL_DARK, radius=18)
    draw.text((124, 766), "Envoyer le code de reinitialisation", font=load_font(20, bold=True), fill=PIL_WHITE)
    draw.text((84, 880), "Etape 2", font=load_font(18, bold=True), fill=PIL_PRIMARY)
    draw.text((84, 918), "Utiliser le code recu", font=load_font(30, bold=True), fill=PIL_TEXT)
    rounded_panel(draw, (84, 972, 520, 1042), PIL_SOFT, radius=18)
    draw.text((108, 994), "Code a 6 chiffres", font=load_font(20), fill=PIL_MUTED)
    rounded_panel(draw, (84, 1060, 520, 1134), PIL_SOFT, radius=18)
    draw.text((108, 1084), "Nouveau mot de passe", font=load_font(20), fill=PIL_MUTED)

    rounded_panel(draw, (620, 174, 1140, 962), PIL_WHITE, radius=34)
    draw.text((662, 218), "Courriel recu", font=load_font(30, bold=True), fill=PIL_TEXT)
    draw.text((662, 266), "SendGrid | Reinitialisation du mot de passe", font=load_font(18), fill=PIL_MUTED)
    rounded_panel(draw, (662, 326, 1098, 486), PIL_SOFT, radius=22)
    draw.text((692, 356), "Code de reinitialisation", font=load_font(20, bold=True), fill=PIL_MUTED)
    rounded_panel(draw, (692, 394, 1068, 462), PIL_DARK, radius=18)
    draw.text((780, 408), "483 216", font=load_font(34, bold=True), fill=PIL_WHITE)
    draw.text((662, 542), "Bonjour Alhousseini,", font=load_font(22, bold=True), fill=PIL_TEXT)
    email_lines = [
        "Nous avons recu une demande de reinitialisation de mot de passe.",
        "Ouvrez l'application Garage Mechanic puis saisissez ce code.",
        "Expiration : 31 mars 2026 15:30",
    ]
    y = 592
    for line in email_lines:
        draw.text((662, y), line, font=load_font(20), fill=PIL_MUTED)
        y += 54
    rounded_panel(draw, (662, 786, 1046, 844), PIL_SUCCESS_BG, radius=16)
    draw.text((692, 804), "Flux valide pour la demo et la production", font=load_font(18, bold=True), fill=PIL_SUCCESS)

    draw.text((74, 1322), "Ecran mobile + courriel transactionnel de reinitialisation", font=load_font(20), fill=PIL_MUTED)
    canvas.save(destination)
    return destination


def render_tutorial_front_visual() -> Path:
    destination = ASSET_DIR / "story_07_tutorial_front_visual.png"
    width, height = 1200, 1480
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    tutorial_preview = fit_image(TUTORIALS_SCREEN, (1020, 1240))
    paste_card(canvas, tutorial_preview, (86, 86, 1114, 1328), radius=34)

    callouts = [
        ((110, 140, 290, 208), "hero"),
        ((106, 420, 376, 490), "recherche"),
        ((118, 584, 360, 648), "filtres"),
        ((706, 500, 1044, 570), "tutoriel a la une"),
        ((734, 950, 1050, 1020), "cartes catalogue"),
    ]
    for bounds, label in callouts:
        rounded_panel(draw, bounds, PIL_PRIMARY, radius=18)
        draw.text((bounds[0] + 18, bounds[1] + 18), label, font=load_font(18, bold=True), fill=PIL_WHITE)

    rounded_panel(draw, (112, 1350, 1088, 1432), PIL_WHITE, radius=22, outline=(226, 232, 240))
    draw.text((144, 1378), "Capture de la nouvelle page tutorielle frontend : hero, recherche, filtres et cartes plus lisibles.", font=load_font(20), fill=PIL_MUTED)

    canvas.save(destination)
    return destination


def PRIMARY_SOFT_TO_PIL() -> tuple[int, int, int]:
    return (254, 226, 226)


def build_cover_visual(visual_paths: list[Path]) -> Path:
    destination = ASSET_DIR / "cover_visual.png"
    width, height = 1400, 900
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    cards = [
        (56, 66, 500, 430),
        (534, 66, 1344, 430),
        (56, 462, 690, 834),
        (722, 462, 1344, 834),
    ]
    selected = [visual_paths[0], visual_paths[1], visual_paths[2], visual_paths[5]]
    sizes = [(412, 332), (780, 332), (602, 340), (590, 340)]

    for path, bounds, size in zip(selected, cards, sizes):
        preview = fit_image(path, size)
        paste_card(canvas, preview, bounds, radius=30)

    rounded_panel(draw, (970, 88, 1290, 166), PIL_PRIMARY, radius=18)
    draw.text((1018, 112), "Sprint 6", font=load_font(26, bold=True), fill=PIL_WHITE)

    canvas.save(destination)
    return destination


def build_visuals() -> dict[str, Path]:
    return {
        "logo_story": render_logo_story_visual(),
        "tutorial_player_story": render_tutorial_playback_visual(),
        "admin_upload_story": render_admin_upload_visual(),
        "admin_routes_story": render_admin_routes_visual(),
        "admin_services_story": render_admin_services_visual(),
        "password_reset_story": render_password_reset_visual(),
        "tutorial_front_story": render_tutorial_front_visual(),
    }


TRACK_LABELS = {
    "mobile": "Mobile / frontend",
    "admin": "Admin web",
    "backend": "Backend / securite",
}


def build_story_statistics() -> dict[str, object]:
    total_stories = len(STORIES)
    total_subtasks = sum(len(story["subtasks"]) for story in STORIES)
    track_counts: dict[str, int] = {key: 0 for key in TRACK_LABELS}
    for story in STORIES:
        track_counts[story["track"]] += 1

    return {
        "total_stories": total_stories,
        "total_subtasks": total_subtasks,
        "total_code_examples": total_stories,
        "track_counts": track_counts,
    }


def add_full_slide_background(slide) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    color: RGBColor = DARK,
    bold: bool = False,
    font_name: str = "Segoe UI",
    align=PP_ALIGN.LEFT,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    paragraph.alignment = align


def add_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str,
    fill_color: RGBColor = WHITE,
    title_color: RGBColor = DARK,
    body_color: RGBColor = SLATE,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.16)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    title_paragraph = frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.text = title
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = title_color
    title_run.font.name = "Segoe UI"

    body_paragraph = frame.add_paragraph()
    body_paragraph.text = body
    body_paragraph.font.size = Pt(13)
    body_paragraph.font.color.rgb = body_color
    body_paragraph.font.name = "Segoe UI"
    body_paragraph.space_before = Pt(8)


def add_points_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    points: list[str],
    title: str = "Points cles",
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    heading = frame.paragraphs[0]
    heading_run = heading.add_run()
    heading_run.text = title
    heading_run.font.size = Pt(18)
    heading_run.font.bold = True
    heading_run.font.color.rgb = DARK
    heading_run.font.name = "Segoe UI"

    for item in points:
        paragraph = frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.level = 0
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = SLATE
        paragraph.font.name = "Segoe UI"
        paragraph.space_before = Pt(5)


def add_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_background(slide)

    stats = build_story_statistics()
    track_counts = stats["track_counts"]

    add_textbox(slide, 0.64, 0.48, 5.8, 0.42, "Vue d'ensemble Sprint 6", 28, bold=True)
    add_textbox(
        slide,
        0.64,
        0.92,
        7.6,
        0.3,
        "Statistiques des stories livrees et volume de sous-taches traitees pendant le sprint.",
        12,
        color=MUTED,
    )

    metric_specs = [
        ("Stories", str(stats["total_stories"]), "stories de demonstration"),
        ("Sous-taches", str(stats["total_subtasks"]), "elements fonctionnels realises"),
        ("Extraits code", str(stats["total_code_examples"]), "snippets commentes dans le deck"),
        ("Focus mobile", str(track_counts["mobile"]), "stories orientees application"),
    ]

    left = 0.64
    for index, (title, value, hint) in enumerate(metric_specs):
        x = left + index * 3.06
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.42),
            Inches(2.74),
            Inches(1.08),
        )
        box.line.fill.background()
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.16)
        frame.margin_top = Inches(0.14)
        p1 = frame.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = SLATE
        r1.font.name = "Segoe UI"
        p2 = frame.add_paragraph()
        p2.text = value
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = DARK
        p2.font.name = "Segoe UI"
        p3 = frame.add_paragraph()
        p3.text = hint
        p3.font.size = Pt(10)
        p3.font.color.rgb = MUTED
        p3.font.name = "Segoe UI"

    add_box(
        slide,
        0.64,
        2.86,
        4.05,
        1.42,
        "Repartition des stories",
        "\n".join(
            [
                f"{TRACK_LABELS['mobile']} : {track_counts['mobile']}",
                f"{TRACK_LABELS['admin']} : {track_counts['admin']}",
                f"{TRACK_LABELS['backend']} : {track_counts['backend']}",
            ]
        ),
    )

    add_points_box(
        slide,
        0.64,
        4.46,
        4.05,
        2.02,
        [
            f"Story {story['number']} - {story['title']} ({len(story['subtasks'])} sous-taches)"
            for story in STORIES
        ],
        title="Carte rapide des stories",
    )

    add_points_box(
        slide,
        4.96,
        2.86,
        7.72,
        3.62,
        [
            "identite visuelle et coherence de marque",
            "lecture video et refonte de l'experience tutorielle",
            "console admin web pour videos et services",
            "routes admin securisees cote backend",
            "reinitialisation de mot de passe par courriel",
        ],
        title="Resume du sprint",
    )


def add_story_slide(prs: Presentation, story: dict[str, object], visuals: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_background(slide)

    number = str(story["number"])
    title = str(story["title"])
    definition = str(story["definition"])
    points = list(story["points"])
    subtasks = list(story["subtasks"])
    code_config = story["code"]

    snippet = extract_snippet(
        code_config["path"],
        code_config["anchor"],
        code_config["before"],
        code_config["after"],
        code_config["max_lines"],
    )
    code_path = ASSET_DIR / f"story_{number}_code.png"
    render_code_panel(code_config["label"], snippet, code_path)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(0.44),
        Inches(1.0),
        Inches(0.42),
    )
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = PRIMARY
    badge_frame = badge.text_frame
    badge_frame.clear()
    badge_paragraph = badge_frame.paragraphs[0]
    badge_run = badge_paragraph.add_run()
    badge_run.text = f"Story {number}"
    badge_run.font.size = Pt(14)
    badge_run.font.bold = True
    badge_run.font.color.rgb = WHITE
    badge_run.font.name = "Segoe UI"
    badge_paragraph.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 1.72, 0.32, 8.2, 0.44, title, 26, bold=True)
    add_textbox(slide, 0.58, 0.96, 8.4, 0.22, "Sprint 6 - demonstration mobile, admin web et backend", 11, color=MUTED)

    slide.shapes.add_picture(
        str(visuals[str(story["visual_key"])]),
        Inches(0.58),
        Inches(1.28),
        width=Inches(5.25),
        height=Inches(5.6),
    )

    add_box(slide, 6.05, 1.28, 6.7, 0.92, "Definition de la story", definition)
    add_points_box(slide, 6.05, 2.34, 6.7, 1.08, points)
    add_points_box(slide, 6.05, 3.56, 6.7, 1.08, subtasks, title="Sous-taches livrees")
    slide.shapes.add_picture(str(code_path), Inches(6.05), Inches(4.82), width=Inches(6.7), height=Inches(1.45))
    add_box(
        slide,
        6.05,
        6.38,
        6.7,
        0.52,
        "Ce que fait ce code dans l'application",
        str(code_config["explanation"]),
        fill_color=PRIMARY_SOFT,
        title_color=PRIMARY,
        body_color=DARK,
    )


def add_cover_slide(prs: Presentation, cover_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_background(slide)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(0.58),
        Inches(1.45),
        Inches(0.44),
    )
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = PRIMARY
    badge_frame = badge.text_frame
    badge_frame.clear()
    badge_paragraph = badge_frame.paragraphs[0]
    badge_run = badge_paragraph.add_run()
    badge_run.text = "Sprint 6"
    badge_run.font.size = Pt(16)
    badge_run.font.bold = True
    badge_run.font.color.rgb = WHITE
    badge_run.font.name = "Segoe UI"
    badge_paragraph.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 0.64, 1.18, 6.1, 0.86, "Presentation de demonstration", 30, bold=True)
    add_textbox(
        slide,
        0.64,
        2.04,
        5.6,
        1.1,
        "Logo, experience tutorielle, console admin, services, routes securisees et reinitialisation par courriel.",
        18,
        color=SLATE,
    )
    add_textbox(slide, 0.64, 3.0, 4.8, 0.24, "Une slide par story, avec visuel et extrait de code commente.", 12, color=MUTED)
    slide.shapes.add_picture(str(cover_path), Inches(6.0), Inches(0.84), width=Inches(6.55), height=Inches(5.92))

    footer = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.64),
        Inches(6.38),
        Inches(12.0),
        Inches(0.5),
    )
    footer.line.fill.background()
    footer.fill.solid()
    footer.fill.fore_color.rgb = WHITE
    footer_frame = footer.text_frame
    footer_frame.clear()
    footer_paragraph = footer_frame.paragraphs[0]
    footer_run = footer_paragraph.add_run()
    footer_run.text = "Projet Mobile Developpement Avancee - deck genere automatiquement depuis le repo"
    footer_run.font.size = Pt(14)
    footer_run.font.color.rgb = DARK
    footer_run.font.name = "Segoe UI"
    footer_paragraph.alignment = PP_ALIGN.CENTER


def save_presentation(prs: Presentation) -> Path:
    candidates = [OUTPUT_PATH]
    candidates.extend(
        OUTPUT_PATH.with_name(f"{OUTPUT_PATH.stem}_v{index}{OUTPUT_PATH.suffix}")
        for index in range(2, 10)
    )

    last_error: PermissionError | None = None
    for candidate in candidates:
        try:
            prs.save(candidate)
            return candidate
        except PermissionError as error:
            last_error = error
            continue

    raise last_error or PermissionError(f"Unable to save presentation: {OUTPUT_PATH}")


def build_presentation() -> Path:
    ensure_asset_dir()
    visuals = build_visuals()
    cover = build_cover_visual(list(visuals.values()))

    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)

    add_cover_slide(prs, cover)
    add_summary_slide(prs)
    for story in STORIES:
        add_story_slide(prs, story, visuals)

    return save_presentation(prs)


if __name__ == "__main__":
    output = build_presentation()
    print(output)
