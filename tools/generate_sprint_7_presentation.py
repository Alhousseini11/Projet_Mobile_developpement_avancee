from __future__ import annotations

import subprocess
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "presentation_assets" / "sprint_7"
OUTPUT_PATH = ROOT / "presentation_sprint_7_demo.pptx"

WIDE_W = 13.333
WIDE_H = 7.5

PRIMARY = RGBColor(220, 38, 38)
PRIMARY_SOFT = RGBColor(254, 226, 226)
DARK = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
SUCCESS = RGBColor(22, 101, 52)
SUCCESS_BG = RGBColor(220, 252, 231)
WARNING = RGBColor(180, 83, 9)
WARNING_BG = RGBColor(255, 237, 213)

PIL_PRIMARY = (220, 38, 38)
PIL_PRIMARY_DARK = (153, 27, 27)
PIL_DARK = (15, 23, 42)
PIL_PANEL = (30, 41, 59)
PIL_SHELL = (241, 245, 249)
PIL_SOFT = (248, 250, 252)
PIL_WHITE = (255, 255, 255)
PIL_MUTED = (100, 116, 139)
PIL_TEXT = (15, 23, 42)
PIL_SUCCESS = (22, 101, 52)
PIL_SUCCESS_BG = (220, 252, 231)
PIL_WARNING = (180, 83, 9)
PIL_WARNING_BG = (255, 237, 213)

MULTISERVICE_SCREEN = ROOT / "presentation_assets" / "sprint_review" / "reservation_multi_mock.png"
TUTORIAL_PLAYER_VISUAL = ROOT / "presentation_assets" / "sprint_6" / "story_02_tutorial_playback_visual.png"

STORIES = [
    {
        "number": "01",
        "title": "Moderation admin des avis",
        "definition": (
            "Un flux de moderation a ete ajoute dans le back-office pour lister les avis clients et "
            "supprimer un commentaire malsain, injurieux ou hors sujet sans impacter la reservation associee."
        ),
        "points": [
            "Vue admin dediee avec client, service, note et commentaire.",
            "Suppression via confirmation UI puis appel DELETE securise.",
            "La liste des avis et les indicateurs du dashboard sont resynchronises apres suppression.",
        ],
        "subtasks": [
            "creation de GET /api/admin/reviews",
            "ajout de DELETE /api/admin/reviews/:reviewId",
            "integration du bouton Supprimer dans la console admin",
        ],
        "problems": [
            "Aucune moderation n'existait cote admin.",
            "Les droits de suppression devaient rester reserves au role ADMIN.",
            "Les compteurs et listes devaient rester coherents apres deletion.",
        ],
        "visual_key": "review_moderation",
        "code": {
            "path": "backend/src/modules/admin/admin.controller.ts",
            "anchor": "export async function deleteAdminReview(req: Request, res: Response) {",
            "before": 2,
            "after": 20,
            "max_lines": 14,
            "label": "backend/src/modules/admin/admin.controller.ts",
            "explanation": (
                "Le controleur valide l'identifiant, verifie que l'avis existe, le supprime puis "
                "renvoie un 204. La moderation reste donc propre et strictement admin."
            ),
        },
    },
    {
        "number": "02",
        "title": "Correction du probleme d'heure en multi-service",
        "definition": (
            "La reservation multi-service n'utilise plus une intersection stricte des creneaux. "
            "Le frontend construit maintenant un plan sequentiel: une heure de debut est choisie, "
            "puis les services suivants s'enchainent sur les prochains slots compatibles."
        ),
        "points": [
            "L'heure de depart redevient visible meme avec plusieurs services.",
            "Les prestations sont planifiees dans l'ordre choisi par l'utilisateur.",
            "Le payload envoye au backend reutilise le meme plan que celui affiche dans l'interface.",
        ],
        "subtasks": [
            "creation du moteur buildReservationSchedulePlans",
            "remplacement de l'intersection brute des slots dans Reservations.vue",
            "reutilisation du plan lors de la creation des reservations",
        ],
        "problems": [
            "L'intersection exacte des horaires renvoyait souvent une liste vide.",
            "L'UI affichait Aucun creneau alors qu'un enchainement reel etait possible.",
            "Le calcul d'affichage et le payload backend devaient rester alignes.",
        ],
        "visual_key": "multiservice_fix",
        "code": {
            "path": "frontend/app/utils/reservationScheduling.ts",
            "ref": "origin/main",
            "anchor": "const nextSlot = sortSlots(slotsByServiceId[serviceId] ?? []).find(slot => {",
            "before": 6,
            "after": 10,
            "max_lines": 14,
            "label": "frontend/app/utils/reservationScheduling.ts",
            "explanation": (
                "Pour chaque service suivant, le moteur cherche le premier slot disponible apres "
                "la fin du precedent. On obtient un vrai plan multi-service au lieu d'une intersection vide."
            ),
        },
    },
    {
        "number": "03",
        "title": "Activation et desactivation des comptes cote admin",
        "definition": (
            "Le back-office permet maintenant de reactiver ou desactiver un compte. "
            "Le statut active est stocke en base, visible dans la table utilisateurs et applique a "
            "l'authentification, au refresh token et aux routes protegees."
        ),
        "points": [
            "Boutons Desactiver / Reactiver directement dans la table admin.",
            "Blocage des connexions et refresh pour un compte inactif.",
            "Protection contre l'auto-desactivation et la desactivation du dernier admin actif.",
        ],
        "subtasks": [
            "ajout du champ active et de la migration Prisma",
            "creation du endpoint PATCH /api/admin/users/:userId/activation",
            "branchage du toggle dans la console admin et dans auth.service",
        ],
        "problems": [
            "Un simple changement visuel ne suffisait pas: le backend devait refuser les sessions inactives.",
            "Il fallait eviter de bloquer la plateforme en desactivant le dernier admin.",
            "Les flux login, refresh et reset devaient rester coherents avec le statut active.",
        ],
        "visual_key": "account_activation",
        "code": {
            "path": "backend/src/modules/admin/admin.controller.ts",
            "ref": "origin/main",
            "anchor": "export async function updateAdminUserActivation(req: Request, res: Response) {",
            "before": 0,
            "after": 54,
            "max_lines": 14,
            "label": "backend/src/modules/admin/admin.controller.ts",
            "explanation": (
                "Le controleur refuse l'auto-desactivation, protege le dernier admin actif, puis "
                "met a jour le champ active. Ce meme statut est ensuite applique par le module d'authentification."
            ),
        },
    },
    {
        "number": "04",
        "title": "Systeme de vues video qualifiees",
        "definition": (
            "Le compteur de vues tutoriels ne repose plus sur un simple clic. "
            "L'application attend une lecture reelle puis le backend dedoublonne par utilisateur avant "
            "d'incrementer le compteur views."
        ),
        "points": [
            "Declenchement cote mobile apres quelques secondes de lecture reelle.",
            "Une seule vue comptee par utilisateur dans la fenetre de qualification.",
            "Mise a jour atomique du compteur dans l'API tutoriels.",
        ],
        "subtasks": [
            "temporisation de lecture dans TutorialVideoPlayer",
            "stockage TutorialView pour memoriser le dernier visionnage",
            "endpoint POST /api/tutorials/:id/views avec garde anti-doublon",
        ],
        "problems": [
            "Un simple tap sur Lire pouvait gonfler artificiellement le compteur.",
            "Il fallait empecher plusieurs increments consecutifs par le meme utilisateur.",
            "Les utilisateurs anonymes devaient ouvrir la video sans compter une vue.",
        ],
        "visual_key": "video_views",
        "code": {
            "path": "backend/src/modules/tutorials/tutorials.controller.ts",
            "ref": "origin/main",
            "anchor": "export async function incrementTutorialViews(req: Request, res: Response) {",
            "before": 0,
            "after": 38,
            "max_lines": 14,
            "label": "backend/src/modules/tutorials/tutorials.controller.ts",
            "explanation": (
                "Le backend charge la trace TutorialView du couple utilisateur-tutoriel, ignore "
                "les vues trop proches puis incremente views dans la meme transaction."
            ),
        },
    },
]


def ensure_asset_dir() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def load_font(size: int, *, mono: bool = False, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    fonts_dir = Path("C:/Windows/Fonts")
    if mono:
        candidates = [
            fonts_dir / ("consolab.ttf" if bold else "consola.ttf"),
            fonts_dir / ("courbd.ttf" if bold else "cour.ttf"),
        ]
    else:
        candidates = [
            fonts_dir / ("segoeuib.ttf" if bold else "segoeui.ttf"),
            fonts_dir / ("arialbd.ttf" if bold else "arial.ttf"),
            fonts_dir / ("calibrib.ttf" if bold else "calibri.ttf"),
        ]

    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)

    return ImageFont.load_default()


def git_show_file(ref: str, relative_path: str) -> list[str]:
    result = subprocess.run(
        ["git", "show", f"{ref}:{relative_path}"],
        cwd=ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
        check=True,
    )
    return result.stdout.splitlines()


def read_file_lines(relative_path: str, ref: str | None = None) -> list[str]:
    if ref:
        return git_show_file(ref, relative_path)

    return (ROOT / relative_path).read_text(encoding="utf-8").splitlines()


def trim_common_indent(lines: list[str]) -> list[str]:
    non_empty = [line for line in lines if line.strip()]
    if not non_empty:
        return lines

    indent = min(len(line) - len(line.lstrip(" ")) for line in non_empty)
    return [line[indent:] if len(line) >= indent else line for line in lines]


def extract_snippet(
    relative_path: str,
    anchor: str,
    before: int,
    after: int,
    max_lines: int,
    ref: str | None = None,
) -> str:
    lines = read_file_lines(relative_path, ref=ref)
    anchor_index = next((index for index, line in enumerate(lines) if anchor in line), None)
    if anchor_index is None:
        raise ValueError(f"Anchor not found in {relative_path}: {anchor}")

    start = max(anchor_index - before, 0)
    end = min(anchor_index + after, len(lines) - 1)
    selected = trim_common_indent(lines[start : end + 1])[:max_lines]
    numbered = [f"{start + offset + 1:>4} | {line.rstrip()}" for offset, line in enumerate(selected)]
    return "\n".join(numbered)


def wrap_code_lines(snippet: str, max_chars: int) -> str:
    wrapped_lines: list[str] = []
    for raw_line in snippet.splitlines():
        if len(raw_line) <= max_chars:
            wrapped_lines.append(raw_line)
            continue

        prefix = raw_line[:7] if "|" in raw_line[:10] else ""
        content = raw_line[7:] if prefix else raw_line
        chunks = textwrap.wrap(
            content,
            width=max_chars - len(prefix),
            replace_whitespace=False,
            drop_whitespace=False,
            break_long_words=False,
            break_on_hyphens=False,
        )
        if not chunks:
            wrapped_lines.append(raw_line)
            continue

        wrapped_lines.append(f"{prefix}{chunks[0]}")
        continuation_prefix = " " * len(prefix)
        for chunk in chunks[1:]:
            wrapped_lines.append(f"{continuation_prefix}{chunk}")
    return "\n".join(wrapped_lines)


def render_code_panel(label: str, snippet: str, explanation: str, destination: Path) -> Path:
    width, height = 1800, 520
    image = Image.new("RGB", (width, height), PIL_DARK)
    draw = ImageDraw.Draw(image)

    draw.rounded_rectangle((0, 0, width - 1, height - 1), radius=34, fill=PIL_DARK)
    draw.rounded_rectangle((0, 0, width - 1, 84), radius=34, fill=PIL_PANEL)
    draw.rounded_rectangle((30, height - 104, width - 30, height - 26), radius=22, fill=(34, 47, 63))

    draw.text((38, 22), label, font=load_font(30, bold=True), fill=PIL_WHITE)
    draw.text((width - 180, 28), "code", font=load_font(20), fill=(148, 163, 184))

    wrapped = wrap_code_lines(snippet, 102)
    top = 118
    line_height = 28
    for index, line in enumerate(wrapped.splitlines()):
        draw.text((40, top + index * line_height), line, font=load_font(20, mono=True), fill=(226, 232, 240))

    explanation_text = textwrap.fill(explanation, width=122)
    draw.text((54, height - 84), explanation_text, font=load_font(18), fill=(226, 232, 240))
    image.save(destination)
    return destination


def fit_image(path: Path, size: tuple[int, int]) -> Image.Image:
    with Image.open(path) as source:
        return ImageOps.fit(source.convert("RGB"), size, method=Image.Resampling.LANCZOS)


def draw_shadow(draw: ImageDraw.ImageDraw, bounds: tuple[int, int, int, int], radius: int = 28) -> None:
    x0, y0, x1, y1 = bounds
    draw.rounded_rectangle((x0 + 12, y0 + 16, x1 + 12, y1 + 16), radius=radius, fill=(221, 228, 235))


def paste_card(base: Image.Image, content: Image.Image, bounds: tuple[int, int, int, int], radius: int = 28) -> None:
    draw = ImageDraw.Draw(base)
    draw_shadow(draw, bounds, radius=radius)
    x0, y0, x1, y1 = bounds
    draw.rounded_rectangle(bounds, radius=radius, fill=PIL_WHITE)
    inner = Image.new("RGB", (x1 - x0 - 28, y1 - y0 - 28), PIL_WHITE)
    inner.paste(content, (0, 0))
    mask = Image.new("L", inner.size, 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, inner.size[0], inner.size[1]), radius=22, fill=255)
    base.paste(inner, (x0 + 14, y0 + 14), mask)


def rounded_panel(
    draw: ImageDraw.ImageDraw,
    bounds: tuple[int, int, int, int],
    fill: tuple[int, int, int],
    radius: int = 26,
    outline: tuple[int, int, int] | None = None,
    width: int = 2,
) -> None:
    if outline:
        draw.rounded_rectangle(bounds, radius=radius, fill=fill, outline=outline, width=width)
    else:
        draw.rounded_rectangle(bounds, radius=radius, fill=fill)


def render_web_shell(draw: ImageDraw.ImageDraw, width: int, height: int, title: str, subtitle: str, chip: str) -> None:
    rounded_panel(draw, (28, 28, width - 28, height - 28), PIL_WHITE, radius=34)
    rounded_panel(draw, (28, 28, width - 28, 150), PIL_DARK, radius=34)
    rounded_panel(draw, (width - 290, 60, width - 82, 126), PIL_PRIMARY, radius=20)
    draw.text((64, 58), "GARAGE MECHANIC", font=load_font(22, bold=True), fill=(248, 250, 252))
    draw.text((64, 96), title, font=load_font(36, bold=True), fill=PIL_WHITE)
    draw.text((64, 188), subtitle, font=load_font(20), fill=PIL_MUTED)
    draw.text((width - 238, 78), chip, font=load_font(22, bold=True), fill=PIL_WHITE)


def PRIMARY_SOFT_TO_PIL() -> tuple[int, int, int]:
    return (254, 226, 226)


def render_multiservice_fix_visual() -> Path:
    destination = ASSET_DIR / "story_02_multiservice_fix_visual.png"
    width, height = 1480, 900
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    mobile_preview = fit_image(MULTISERVICE_SCREEN, (520, 760))
    paste_card(canvas, mobile_preview, (48, 64, 604, 852), radius=34)

    rounded_panel(draw, (648, 64, 1420, 264), PIL_DARK, radius=30)
    draw.text((684, 104), "Probleme constate", font=load_font(30, bold=True), fill=PIL_WHITE)
    draw.text((684, 154), "Avec plusieurs services, l'intersection stricte des slots pouvait\nsupprimer toute heure visible dans l'interface.", font=load_font(21), fill=(203, 213, 225))

    rounded_panel(draw, (648, 300, 1420, 602), PIL_WHITE, radius=30, outline=(226, 232, 240))
    draw.text((684, 338), "Correction appliquee", font=load_font(30, bold=True), fill=PIL_TEXT)
    steps = [
        ("10:00", "Batterie", "heure de depart affichee"),
        ("10:30", "Diagnostic", "slot suivant compatible"),
        ("11:00", "Fin du plan", "sequence multi-service coherente"),
    ]
    y = 398
    for time_label, title, copy in steps:
        rounded_panel(draw, (684, y, 1380, y + 58), PIL_SOFT, radius=18)
        rounded_panel(draw, (704, y + 12, 808, y + 46), PIL_PRIMARY, radius=12)
        draw.text((732, y + 18), time_label, font=load_font(18, bold=True), fill=PIL_WHITE)
        draw.text((838, y + 12), title, font=load_font(22, bold=True), fill=PIL_TEXT)
        draw.text((1038, y + 14), copy, font=load_font(18), fill=PIL_MUTED)
        y += 74

    rounded_panel(draw, (648, 638, 1420, 852), PIL_WARNING_BG, radius=30)
    draw.text((684, 674), "Effet cote utilisateur", font=load_font(30, bold=True), fill=PIL_WARNING)
    bullets = [
        "le choix d'heure reapparait en multi-service",
        "le recapitulatif conserve une heure de debut claire",
        "la creation des reservations reemploie le meme plan",
    ]
    y = 728
    for bullet in bullets:
        rounded_panel(draw, (686, y + 8, 704, y + 26), PIL_WARNING, radius=8)
        draw.text((724, y), bullet, font=load_font(20), fill=PIL_TEXT)
        y += 42

    canvas.save(destination)
    return destination


def render_review_moderation_visual() -> Path:
    destination = ASSET_DIR / "story_01_review_moderation_visual.png"
    width, height = 1480, 900
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    render_web_shell(
        draw,
        width,
        height,
        "Console admin | Moderation des avis",
        "Suppression des commentaires malsains, injurieux ou hors sujet depuis le back-office.",
        "Story 01",
    )

    rounded_panel(draw, (60, 252, 940, 820), PIL_WHITE, radius=30)
    draw.text((96, 292), "Moderation des avis", font=load_font(34, bold=True), fill=PIL_TEXT)
    draw.text((96, 336), "Liste admin rechargee apres suppression", font=load_font(20), fill=PIL_MUTED)

    cards = [
        ("Amadou Keita", "Diagnostic", "1/5", "Commentaire agressif a moderer immediatement.", True),
        ("Sophie Martin", "Vidange", "5/5", "Intervention rapide et atelier ponctuel.", False),
        ("Karim Diallo", "Freins", "2/5", "Message insultant hors du cadre du service.", True),
    ]
    y = 390
    for customer, service, rating, copy, flagged in cards:
        rounded_panel(draw, (96, y, 904, y + 112), PIL_SOFT if not flagged else (255, 248, 248), radius=22, outline=(226, 232, 240))
        draw.text((124, y + 18), customer, font=load_font(22, bold=True), fill=PIL_TEXT)
        draw.text((124, y + 50), f"{service} | {rating}", font=load_font(18), fill=PIL_MUTED)
        draw.text((124, y + 78), copy, font=load_font(18), fill=PIL_TEXT if flagged else PIL_MUTED)
        button_fill = PIL_PRIMARY if flagged else (226, 232, 240)
        button_text = PIL_WHITE if flagged else PIL_TEXT
        rounded_panel(draw, (742, y + 30, 876, y + 76), button_fill, radius=16)
        draw.text((774, y + 42), "Supprimer", font=load_font(18, bold=True), fill=button_text)
        y += 134

    rounded_panel(draw, (980, 252, 1418, 820), PIL_DARK, radius=30)
    draw.text((1016, 292), "Flux de moderation", font=load_font(30, bold=True), fill=PIL_WHITE)
    callouts = [
        "confirmation avant deletion",
        "DELETE /api/admin/reviews/:reviewId",
        "refresh du dashboard apres succes",
    ]
    y = 364
    for item in callouts:
        rounded_panel(draw, (1016, y, 1382, y + 74), PIL_PANEL, radius=18)
        draw.text((1040, y + 24), item, font=load_font(20), fill=(226, 232, 240))
        y += 96

    rounded_panel(draw, (1016, 672, 1382, 742), PIL_SUCCESS_BG, radius=18)
    draw.text((1044, 694), "Avis supprime avec succes.", font=load_font(20, bold=True), fill=PIL_SUCCESS)
    draw.text((1016, 774), "L'admin garde la main sur la qualite du contenu visible dans la plateforme.", font=load_font(19), fill=(203, 213, 225))

    canvas.save(destination)
    return destination


def render_account_activation_visual() -> Path:
    destination = ASSET_DIR / "story_03_account_activation_visual.png"
    width, height = 1480, 900
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    render_web_shell(
        draw,
        width,
        height,
        "Console admin | Utilisateurs",
        "Activation et desactivation des acces sans quitter la table des comptes.",
        "Story 03",
    )

    rounded_panel(draw, (60, 250, 1418, 820), PIL_WHITE, radius=30)
    draw.text((96, 290), "Gestion des comptes", font=load_font(34, bold=True), fill=PIL_TEXT)
    headers = [
        (96, "Utilisateur"),
        (434, "Role"),
        (606, "Statut"),
        (810, "Activite"),
        (1192, "Action"),
    ]
    for x, label in headers:
        draw.text((x, 340), label, font=load_font(18, bold=True), fill=PIL_MUTED)

    rows = [
        ("Adam Admin", "ADMIN", "Actif", "12 vehicules | 8 reservations", "Session actuelle"),
        ("Moussa Client", "USER", "Actif", "2 vehicules | 3 reservations", "Desactiver"),
        ("Awa Diallo", "USER", "Desactive", "1 vehicule | 1 reservation", "Reactiver"),
    ]
    y = 380
    for full_name, role, status, activity, action in rows:
        rounded_panel(draw, (90, y, 1382, y + 96), PIL_SOFT, radius=20)
        draw.text((116, y + 18), full_name, font=load_font(22, bold=True), fill=PIL_TEXT)
        draw.text((116, y + 52), "contact@example.com", font=load_font(18), fill=PIL_MUTED)
        draw.text((446, y + 34), role, font=load_font(20, bold=True), fill=PIL_TEXT)

        badge_fill = PIL_SUCCESS_BG if status == "Actif" else PRIMARY_SOFT_TO_PIL()
        badge_text = PIL_SUCCESS if status == "Actif" else PIL_PRIMARY_DARK
        rounded_panel(draw, (596, y + 24, 742, y + 68), badge_fill, radius=16)
        draw.text((630, y + 36), status, font=load_font(18, bold=True), fill=badge_text)

        draw.text((810, y + 34), activity, font=load_font(18), fill=PIL_TEXT)

        if action == "Session actuelle":
            rounded_panel(draw, (1180, y + 24, 1354, y + 68), (226, 232, 240), radius=16)
            draw.text((1200, y + 36), action, font=load_font(16, bold=True), fill=PIL_TEXT)
        elif action == "Desactiver":
            rounded_panel(draw, (1194, y + 20, 1358, y + 72), PIL_PRIMARY_DARK, radius=18)
            draw.text((1224, y + 36), action, font=load_font(18, bold=True), fill=PIL_WHITE)
        else:
            rounded_panel(draw, (1202, y + 20, 1358, y + 72), (226, 232, 240), radius=18)
            draw.text((1234, y + 36), action, font=load_font(18, bold=True), fill=PIL_TEXT)

        y += 118

    rounded_panel(draw, (60, 840, 1418, 864), PIL_SOFT, radius=12)
    draw.text((88, 844), "Garde-fous backend: impossible de desactiver son propre compte ou le dernier admin actif.", font=load_font(18), fill=PIL_MUTED)

    canvas.save(destination)
    return destination


def render_video_views_visual() -> Path:
    destination = ASSET_DIR / "story_04_video_views_visual.png"
    width, height = 1480, 900
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    player_preview = fit_image(TUTORIAL_PLAYER_VISUAL, (620, 760))
    paste_card(canvas, player_preview, (48, 64, 704, 852), radius=34)

    rounded_panel(draw, (748, 64, 1420, 270), PIL_DARK, radius=30)
    draw.text((784, 102), "Regle de qualification", font=load_font(30, bold=True), fill=PIL_WHITE)
    draw.text((784, 150), "La vue n'est pas comptee au simple clic.\nLe lecteur attend une lecture reelle avant l'appel API.", font=load_font(21), fill=(203, 213, 225))

    rounded_panel(draw, (748, 308, 1420, 602), PIL_WHITE, radius=30, outline=(226, 232, 240))
    draw.text((784, 346), "Pipeline de comptage", font=load_font(30, bold=True), fill=PIL_TEXT)
    pipeline = [
        "1. le mobile attend 5 secondes de lecture",
        "2. le backend recharge TutorialView(userId, tutorialId)",
        "3. views s'incremente une seule fois dans la transaction",
    ]
    y = 404
    for item in pipeline:
        rounded_panel(draw, (786, y + 8, 804, y + 26), PIL_PRIMARY, radius=8)
        draw.text((824, y), item, font=load_font(19), fill=PIL_TEXT)
        y += 52

    rounded_panel(draw, (748, 640, 1420, 852), PIL_SUCCESS_BG, radius=30)
    draw.text((784, 678), "Resultat pour la demo", font=load_font(30, bold=True), fill=PIL_SUCCESS)
    draw.text((784, 728), "Le compteur reste credible: pas de gonflement artificiel, pas de double vue immediate, et un suivi utilisateur propre.", font=load_font(21), fill=PIL_TEXT)

    canvas.save(destination)
    return destination


def build_cover_visual(visual_paths: list[Path]) -> Path:
    destination = ASSET_DIR / "cover_visual.png"
    width, height = 1400, 900
    canvas = Image.new("RGB", (width, height), PIL_SHELL)
    draw = ImageDraw.Draw(canvas)

    cards = [
        (52, 70, 492, 424),
        (520, 70, 1348, 424),
        (52, 454, 680, 838),
        (710, 454, 1348, 838),
    ]
    sizes = [(408, 322), (796, 322), (596, 352), (606, 352)]

    for path, bounds, size in zip(visual_paths, cards, sizes):
        preview = fit_image(path, size)
        paste_card(canvas, preview, bounds, radius=30)

    rounded_panel(draw, (1030, 94, 1300, 166), PIL_PRIMARY, radius=18)
    draw.text((1090, 116), "Sprint 7", font=load_font(26, bold=True), fill=PIL_WHITE)
    canvas.save(destination)
    return destination


def build_visuals() -> dict[str, Path]:
    return {
        "review_moderation": render_review_moderation_visual(),
        "multiservice_fix": render_multiservice_fix_visual(),
        "account_activation": render_account_activation_visual(),
        "video_views": render_video_views_visual(),
    }


def build_story_statistics() -> dict[str, int]:
    return {
        "total_stories": len(STORIES),
        "total_subtasks": sum(len(story["subtasks"]) for story in STORIES),
        "total_code_examples": len(STORIES),
        "total_modules": 3,
    }


def add_full_slide_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int,
    *,
    color: RGBColor = DARK,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    if align is not None:
        paragraph.alignment = align


def add_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str,
    *,
    fill_color: RGBColor = WHITE,
    title_color: RGBColor = DARK,
    body_color: RGBColor = SLATE,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.16)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    heading = frame.paragraphs[0]
    heading_run = heading.add_run()
    heading_run.text = title
    heading_run.font.size = Pt(17)
    heading_run.font.bold = True
    heading_run.font.color.rgb = title_color
    heading_run.font.name = "Segoe UI"

    paragraph = frame.add_paragraph()
    paragraph.text = body
    paragraph.font.size = Pt(12)
    paragraph.font.color.rgb = body_color
    paragraph.font.name = "Segoe UI"
    paragraph.space_before = Pt(6)


def add_bullet_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    items: list[str],
    *,
    fill_color: RGBColor = WHITE,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    heading = frame.paragraphs[0]
    heading_run = heading.add_run()
    heading_run.text = title
    heading_run.font.size = Pt(17)
    heading_run.font.bold = True
    heading_run.font.color.rgb = DARK
    heading_run.font.name = "Segoe UI"

    for item in items:
        paragraph = frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.level = 0
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = SLATE
        paragraph.font.name = "Segoe UI"
        paragraph.space_before = Pt(4)


def add_metric_card(slide, left: float, top: float, title: str, value: str, hint: str) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(2.82),
        Inches(1.02),
    )
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.16)
    frame.margin_top = Inches(0.12)

    p1 = frame.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = SLATE
    r1.font.name = "Segoe UI"

    p2 = frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = DARK
    p2.font.name = "Segoe UI"

    p3 = frame.add_paragraph()
    p3.text = hint
    p3.font.size = Pt(10)
    p3.font.color.rgb = MUTED
    p3.font.name = "Segoe UI"


def add_intro_slide(prs: Presentation, cover_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_background(slide)

    stats = build_story_statistics()

    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.62),
        Inches(0.54),
        Inches(1.52),
        Inches(0.44),
    )
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = PRIMARY
    chip_frame = chip.text_frame
    chip_frame.clear()
    chip_p = chip_frame.paragraphs[0]
    chip_run = chip_p.add_run()
    chip_run.text = "Sprint 7"
    chip_run.font.size = Pt(16)
    chip_run.font.bold = True
    chip_run.font.color.rgb = WHITE
    chip_run.font.name = "Segoe UI"
    chip_p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 0.64, 1.18, 5.7, 0.82, "Presentation de demonstration", 30, bold=True)
    add_textbox(
        slide,
        0.64,
        2.0,
        5.5,
        0.82,
        "Moderation des avis, fiabilite multi-service, gestion des comptes et vues video qualifiees.",
        18,
        color=SLATE,
    )
    add_textbox(
        slide,
        0.64,
        2.78,
        5.3,
        0.24,
        "Perimetre technique : mobile NativeScript-Vue, admin web TypeScript, backend Node/Express/Prisma.",
        11,
        color=MUTED,
    )

    slide.shapes.add_picture(str(cover_path), Inches(6.35), Inches(0.78), width=Inches(6.3), height=Inches(3.7))

    add_metric_card(slide, 0.64, 3.16, "Stories", str(stats["total_stories"]), "stories de demo")
    add_metric_card(slide, 3.64, 3.16, "Sous-taches", str(stats["total_subtasks"]), "elements fonctionnels")
    add_metric_card(slide, 0.64, 4.34, "Extraits code", str(stats["total_code_examples"]), "snippets commentes")
    add_metric_card(slide, 3.64, 4.34, "Couches", str(stats["total_modules"]), "mobile, admin, backend")

    story_cards = [
        (0.64, 5.56, STORIES[0]),
        (3.82, 5.56, STORIES[1]),
        (7.02, 5.56, STORIES[2]),
        (10.2, 5.56, STORIES[3]),
    ]

    for left, top, story in story_cards:
        add_bullet_box(
            slide,
            left,
            top,
            2.5,
            1.44,
            f"Story {story['number']} - {story['title']}",
            list(story["subtasks"]),
            fill_color=WHITE,
        )


def add_story_slide(prs: Presentation, story: dict[str, object], visuals: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_background(slide)

    number = str(story["number"])
    title = str(story["title"])
    code_config = story["code"]

    snippet = extract_snippet(
        code_config["path"],
        code_config["anchor"],
        code_config["before"],
        code_config["after"],
        code_config["max_lines"],
        ref=code_config.get("ref"),
    )
    code_path = ASSET_DIR / f"story_{number}_code.png"
    render_code_panel(code_config["label"], snippet, code_config["explanation"], code_path)

    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.58),
        Inches(0.4),
        Inches(1.02),
        Inches(0.42),
    )
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = PRIMARY
    chip_frame = chip.text_frame
    chip_frame.clear()
    chip_p = chip_frame.paragraphs[0]
    chip_run = chip_p.add_run()
    chip_run.text = f"Story {number}"
    chip_run.font.size = Pt(14)
    chip_run.font.bold = True
    chip_run.font.color.rgb = WHITE
    chip_run.font.name = "Segoe UI"
    chip_p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 1.72, 0.32, 8.8, 0.44, title, 26, bold=True)
    add_textbox(
        slide,
        0.58,
        0.94,
        8.9,
        0.22,
        "Sprint 7 - demonstration mobile, admin web et backend",
        11,
        color=MUTED,
    )

    slide.shapes.add_picture(
        str(visuals[str(story["visual_key"])]),
        Inches(0.58),
        Inches(1.24),
        width=Inches(12.15),
        height=Inches(2.52),
    )

    add_box(slide, 0.58, 3.98, 5.86, 1.0, "Definition de la story", str(story["definition"]))
    add_bullet_box(slide, 6.54, 3.98, 6.19, 1.0, "Points cles", list(story["points"]))
    add_bullet_box(slide, 0.58, 5.14, 5.86, 0.82, "Sous-taches livrees", list(story["subtasks"]))
    add_bullet_box(slide, 6.54, 5.14, 6.19, 0.82, "Problemes resolus", list(story["problems"]), fill_color=WARNING_BG)

    slide.shapes.add_picture(
        str(code_path),
        Inches(0.58),
        Inches(6.12),
        width=Inches(12.15),
        height=Inches(1.04),
    )


def save_presentation(prs: Presentation) -> Path:
    candidates = [OUTPUT_PATH]
    candidates.extend(
        OUTPUT_PATH.with_name(f"{OUTPUT_PATH.stem}_v{index}{OUTPUT_PATH.suffix}")
        for index in range(2, 10)
    )

    last_error: PermissionError | None = None
    for candidate in candidates:
        try:
            prs.save(candidate)
            return candidate
        except PermissionError as error:
            last_error = error
            continue

    raise last_error or PermissionError(f"Unable to save presentation: {OUTPUT_PATH}")


def build_presentation() -> Path:
    ensure_asset_dir()
    visuals = build_visuals()
    cover = build_cover_visual(
        [
            visuals["review_moderation"],
            visuals["multiservice_fix"],
            visuals["account_activation"],
            visuals["video_views"],
        ]
    )

    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)

    add_intro_slide(prs, cover)
    for story in STORIES:
        add_story_slide(prs, story, visuals)

    return save_presentation(prs)


if __name__ == "__main__":
    output = build_presentation()
    print(output)
