from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "presentation_sprint_8_review.pptx"
LOGO_PATH = ROOT / "frontend" / "app" / "assets" / "Logo-centré.png"

WIDE_W = 13.333
WIDE_H = 7.5

PRIMARY = RGBColor(220, 38, 38)
PRIMARY_DARK = RGBColor(153, 27, 27)
NAVY = RGBColor(17, 24, 39)
NAVY_LIGHT = RGBColor(31, 41, 55)
SLATE = RGBColor(71, 85, 105)
TEXT = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(241, 245, 249)
LIGHTER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
SUCCESS = RGBColor(21, 128, 61)
SUCCESS_BG = RGBColor(220, 252, 231)
READY = RGBColor(180, 83, 9)
READY_BG = RGBColor(255, 237, 213)
INFO = RGBColor(29, 78, 216)
INFO_BG = RGBColor(219, 234, 254)

STATUS_STYLES = {
    "Mergée sur main": (SUCCESS, SUCCESS_BG),
    "Branche prête à merger": (READY, READY_BG),
    "Complément de story": (INFO, INFO_BG),
}

SUMMARY_CARDS = [
    ("10 stories", "socle prod, qualité, maintenabilité et runtime mobile"),
    ("72 tests backend", "validés localement avec la suite complète verte"),
    ("8 migrations Prisma", "vérifiées, appliquées et contrôlées par le preflight"),
    ("3 environnements frontend", "local, VPS partagé, prod explicitement séparés"),
]

STORIES = [
    {
        "number": "01",
        "title": "Corriger le drift Prisma / base de données",
        "status": "Mergée sur main",
        "problem": [
            "Le backend pouvait démarrer sur une base incohérente avec le schéma Prisma.",
            "Le drift n'était détecté qu'après coup, souvent pendant les tests ou le déploiement.",
        ],
        "subtasks": [
            "Construire un snapshot de compatibilité base <-> schéma Prisma.",
            "Détecter l'absence de _prisma_migrations, les migrations manquantes, en échec ou inconnues.",
            "Comparer les tables et colonnes réelles avec les modèles Prisma attendus.",
            "Bloquer le démarrage si la base est incompatible.",
            "Documenter le workflow Prisma et la conduite à tenir en cas de drift.",
        ],
        "delivered": [
            "Ajout de databaseReadiness.ts et schemaManifest.ts.",
            "Nouveau preflight Prisma au démarrage du serveur.",
            "Commande réutilisable prisma:check pour CI, dev et prod.",
            "Workflow Docker aligné sur prisma migrate deploy avant node dist/server.js.",
        ],
        "impact": [
            "Échec explicite au démarrage au lieu d'un backend semi-fonctionnel.",
            "Déploiements plus sûrs et diagnostic plus rapide.",
            "Base de données traitée comme un artefact versionné, pas comme un état implicite.",
        ],
        "files": [
            "backend/src/data/prisma/databaseReadiness.ts",
            "backend/src/data/prisma/schemaManifest.ts",
            "backend/src/server.ts",
            "backend/PRISMA_WORKFLOW.md",
        ],
        "evidence": [
            "Erreur produite: missing \"_prisma_migrations\" table; the database was not initialized with Prisma Migrate",
            "Preuve de recette: prisma:status, prisma:check et prisma:migrate:deploy passent sur le VPS.",
        ],
    },
    {
        "number": "02",
        "title": "Retirer les secrets du Dockerfile",
        "status": "Mergée sur main",
        "problem": [
            "L'image backend embarquait trop de suppositions sur la base et le runtime.",
            "Le risque principal était de figer de mauvaises valeurs ou d'exposer des secrets au build.",
        ],
        "subtasks": [
            "Garder uniquement une DATABASE_URL factice pour le build Prisma.",
            "Injecter la vraie configuration uniquement au runtime via Docker Compose.",
            "Conserver prisma migrate deploy dans le conteneur de prod.",
            "Documenter clairement quelles variables viennent du runtime et non de l'image.",
        ],
        "delivered": [
            "Dockerfile multi-stage propre avec placeholder build-only.",
            "Suppression des identifiants runtime du Dockerfile.",
            "Commande d'entrée: vérification DATABASE_URL puis prisma migrate deploy puis démarrage.",
        ],
        "impact": [
            "Image backend réutilisable entre environnements.",
            "Séparation nette build / runtime.",
            "Réduction du risque de fuite ou de mauvaise configuration en production.",
        ],
        "files": [
            "backend/Dockerfile",
            "docker-compose.prod.yml",
            "DEPLOYMENT.md",
        ],
        "evidence": [
            "La vraie DATABASE_URL est injectée par .env.prod au runtime.",
            "Le backend refuse maintenant de démarrer si la variable manque.",
        ],
    },
    {
        "number": "03",
        "title": "Durcir env.ts pour la production",
        "status": "Mergée sur main",
        "problem": [
            "JWT_SECRET avait un fallback dangereux et DEMO_MODE pouvait rester actif par défaut.",
            "Les erreurs d'environnement n'étaient pas suffisamment centralisées ni explicites.",
        ],
        "subtasks": [
            "Ajouter requireEnv(name) et requireEnvInProduction(name).",
            "Exiger DATABASE_URL et JWT_SECRET en production.",
            "Désactiver DEMO_MODE par défaut.",
            "Valider les variables critiques et les origines CORS au démarrage.",
            "Mettre à jour .env.example, .env.prod.example, README et DEPLOYMENT.",
        ],
        "delivered": [
            "Validation runtime centralisée dans env.ts.",
            "Message d'erreur lisible et bloquant si une variable critique manque.",
            "Exemples d'environnement alignés avec les garde-fous de production.",
        ],
        "impact": [
            "Le backend échoue tôt, proprement et avec un message actionnable.",
            "Le compte de démo n'est plus exposé par erreur en prod.",
            "La recette VPS a confirmé la valeur: l'absence de JWT_SECRET a été détectée immédiatement.",
        ],
        "files": [
            "backend/src/config/env.ts",
            "backend/.env.example",
            "backend/.env.prod.example",
            "README.md",
            "DEPLOYMENT.md",
        ],
        "evidence": [
            "Exemple réel: Missing required environment variable \"JWT_SECRET\" in production.",
            "Le diagnostic a évité un faux problème réseau ou Prisma pendant le déploiement.",
        ],
    },
    {
        "number": "04",
        "title": "Remplacer le fallback d'IP VPS dans le frontend mobile",
        "status": "Mergée sur main",
        "problem": [
            "L'application mobile pouvait tomber sur la prod par défaut à cause d'une IP codée en dur.",
            "Les contextes local, VPS partagé et prod n'étaient pas assez explicites.",
        ],
        "subtasks": [
            "Lire NS_API_BASE_URL au lieu d'un fallback implicite.",
            "Normaliser l'URL et rejeter les valeurs non absolues ou non http/https.",
            "Créer une convention d'environnement explicite: .env.local, .env.shared-vps, .env.prod.",
            "Fournir des scripts dédiés android:local, android:shared-vps, build:android:prod.",
            "Documenter les usages dans README et ARCHITECTURE.",
        ],
        "delivered": [
            "api.ts durci avec erreur claire si l'URL n'est pas configurée.",
            "Scripts run-with-env.js pour charger le bon .env côté NativeScript.",
            "Exemples .env.* propres pour local, VPS et prod.",
        ],
        "impact": [
            "Plus de bascule accidentelle vers la prod.",
            "Débogage facilité: l'erreur d'environnement est explicite dès le démarrage.",
            "Les développeurs choisissent désormais consciemment leur cible API.",
        ],
        "files": [
            "frontend/app/utils/api.ts",
            "frontend/scripts/run-with-env.js",
            "frontend/package.json",
            "frontend/.env.local.example",
            "frontend/.env.shared-vps.example",
            "frontend/.env.prod.example",
        ],
        "evidence": [
            "Message utilisateur: [API_CONFIG] NS_API_BASE_URL is not configured.",
            "Commandes dédiées: npm run android:local, npm run android:shared-vps, build:android:prod.",
        ],
    },
    {
        "number": "05",
        "title": "Nettoyer et fiabiliser la CI GitHub Actions",
        "status": "Mergée sur main",
        "problem": [
            "La CI avait des versions Node désalignées et des étapes redondantes.",
            "Le coverage backend échouait sur Node 20 avec les anciens flags de test coverage.",
        ],
        "subtasks": [
            "Aligner toute la CI sur Node 20.",
            "Faire générer et valider Prisma avant les tests.",
            "Appliquer les migrations sur une base CI propre puis lancer prisma:check.",
            "Séparer clairement backend tests, backend coverage et frontend type-check via une matrix lisible.",
            "Remplacer la couverture incompatible par c8.",
        ],
        "delivered": [
            "Workflow CI plus court et plus compréhensible.",
            "Script test:coverage:ci stabilisé avec c8.",
            "Étapes ordonnées: install -> prisma generate -> validate -> migrate -> check -> test.",
        ],
        "impact": [
            "La CI détecte les problèmes de schéma et de compatibilité avant les tests fonctionnels.",
            "Moins de doublons entre jobs et un signal plus fiable.",
            "Support natif de la couverture backend sur Node 20.",
        ],
        "files": [
            ".github/workflows/ci.yml",
            "backend/package.json",
        ],
        "evidence": [
            "Pipeline frontend: npm run type-check.",
            "Pipeline backend: npm test et npm run test:coverage:ci.",
        ],
    },
    {
        "number": "06",
        "title": "Refactoriser le module home pour le rendre testable",
        "status": "Mergée sur main",
        "problem": [
            "La logique du home feed mélangeait HTTP, métier, accès données et formatage.",
            "Cette zone était sensible et difficile à tester proprement.",
        ],
        "subtasks": [
            "Extraire la logique métier dans un vrai home.service.",
            "Isoler l'accès Prisma dans home.repository.",
            "Conserver un contrôleur HTTP fin et compatible avec l'API existante.",
            "Créer des helpers métiers purs pour l'affichage, les rendez-vous et les rappels.",
            "Ajouter des tests unitaires ciblés sur les cas succès, vide et erreur.",
        ],
        "delivered": [
            "Découpage clair home.controller / home.service / home.repository.",
            "Injection des dépendances pour tester sans Prisma réel.",
            "Fallback métier propre quand les dépendances échouent.",
        ],
        "impact": [
            "Le home feed est devenu lisible, isolé et maintenable.",
            "Les règles de formatage sont testables sans environnement HTTP.",
            "Le contrôleur ne fait plus de travail métier caché.",
        ],
        "files": [
            "backend/src/modules/home/home.controller.ts",
            "backend/src/modules/home/home.service.ts",
            "backend/src/modules/home/home.repository.ts",
            "backend/tests/homeService.test.ts",
        ],
        "evidence": [
            "Cas testés: feed enrichi, feed vide valide, fallback si dépendance en erreur.",
        ],
    },
    {
        "number": "07",
        "title": "Ajouter des tests vraiment utiles sur les routes critiques",
        "status": "Mergée sur main",
        "problem": [
            "Le besoin n'était pas un pourcentage abstrait, mais la sécurisation des parcours métier critiques.",
            "Les zones prioritaires étaient auth, home, reservations, profile et password reset.",
        ],
        "subtasks": [
            "Monter un setup d'intégration avec base PostgreSQL temporaire et migrations appliquées à chaud.",
            "Couvrir les cas succès, erreur métier et autorisation.",
            "Tester les routes publiques, les routes protégées et les comportements de sécurité.",
            "Ajouter des tests unitaires ciblés quand un service est mieux isolé qu'en intégration.",
        ],
        "delivered": [
            "Suite backend complète: 72 tests verts.",
            "Scénarios auth/login/refresh/profile complets.",
            "Password reset sans enumeration de comptes et invalidation des codes réutilisés.",
            "Health endpoint, home feed sans données, erreurs 401/403/404/409 réalistes.",
        ],
        "impact": [
            "Les parcours à plus fort risque sont couverts par des comportements observables.",
            "La refacto et la CI peuvent évoluer avec un filet de sécurité réel.",
            "Les régressions de sécurité deviennent visibles rapidement.",
        ],
        "files": [
            "backend/tests/httpIntegration.test.ts",
            "backend/tests/homeService.test.ts",
            "backend/tests/health.test.ts",
            "backend/tests/controllerInternals.test.ts",
        ],
        "evidence": [
            "Validation locale exécutée: 72 tests, 72 pass, 0 fail.",
            "Exemples couverts: email déjà utilisé -> 409, notifications sans auth -> 401, home invalide mais stable -> 200.",
        ],
    },
    {
        "number": "08",
        "title": "Ajouter observabilité et healthcheck utile",
        "status": "Mergée sur main",
        "problem": [
            "Le backend était difficile à diagnostiquer en cas de panne ou de lenteur.",
            "Le healthcheck ne donnait pas assez d'information opérationnelle.",
        ],
        "subtasks": [
            "Ajouter un request id et le propager dans la réponse.",
            "Logger méthode, URL, code HTTP, durée, taille de réponse, userId et IP.",
            "Rendre les erreurs Prisma plus lisibles via une description structurée.",
            "Créer un /health avec statut app, statut DB, timestamp et environnement.",
        ],
        "delivered": [
            "Middleware requestLogger avec x-request-id et durationMs.",
            "Healthcheck applicatif + sonde base de données.",
            "Logs structurés exploitables pour Docker, CI et prod.",
        ],
        "impact": [
            "Diagnostic plus rapide des erreurs backend et des lenteurs.",
            "Healthcheck réutilisable par curl, Docker healthcheck, CI et supervision simple.",
            "Les erreurs Prisma sont mieux comprises sans exposer de données sensibles au client.",
        ],
        "files": [
            "backend/src/core/http/middleware/requestLogger.ts",
            "backend/src/core/http/health.ts",
            "backend/src/data/prisma/prismaError.ts",
            "backend/tests/health.test.ts",
        ],
        "evidence": [
            "Exemple health: ok/status/environment/timestamp/checks.app/checks.db.",
            "Exemple log: method, url, statusCode, requestId, durationMs, userId, responseBytes.",
        ],
    },
    {
        "number": "09",
        "title": "Durcir la sécurité HTTP minimale du backend",
        "status": "Branche prête à merger",
        "problem": [
            "La couche Express restait permissive: pas de helmet, CORS trop ouvert et erreurs génériques trop bavardes.",
            "Les routes auth et reset de mot de passe n'étaient pas freinées côté HTTP.",
        ],
        "subtasks": [
            "Ajouter helmet avec configuration compatible mobile.",
            "Introduire un CORS explicite par environnement et une allowlist navigateur.",
            "Limiter la taille des payloads JSON et formulaires.",
            "Ajouter un rate limit minimal sur register, login, forgot-password, reset-password et le formulaire HTML.",
            "Masquer les erreurs génériques en production tout en préservant les AppError métier.",
        ],
        "delivered": [
            "Nouveau module security.ts pour CORS et rate limit.",
            "httpErrors.ts pour normaliser les statuts et les messages publics.",
            "Mise à jour de createHttpApp, errorHandler, env.ts, README et DEPLOYMENT.",
            "Tests dédiés httpSecurity.test.ts.",
        ],
        "impact": [
            "Réduction de l'exposition aux abus de base et aux erreurs trop verbeuses.",
            "Configuration explicite pour un backend derrière Nginx ou load balancer.",
            "Lot prêt à merger sans refonte lourde de l'architecture.",
        ],
        "files": [
            "backend/src/core/http/security.ts",
            "backend/src/core/http/httpErrors.ts",
            "backend/src/core/http/createHttpApp.ts",
            "backend/tests/httpSecurity.test.ts",
            "branche: feat/backend-http-security-hardening",
        ],
        "evidence": [
            "4 routes auth protégées par rate limit.",
            "Erreurs génériques masquées en production: Internal error au lieu d'un stack ou message brut.",
        ],
    },
    {
        "number": "10",
        "title": "Optimiser la fiabilité perçue des appels API frontend",
        "status": "Branche prête à merger",
        "problem": [
            "Sur Android NativeScript, le bon .env n'était pas toujours injecté dans le bundle runtime.",
            "Résultat visible côté utilisateur: faux timeouts, crashs au démarrage ou app connectée à la mauvaise cible.",
        ],
        "subtasks": [
            "Injecter NS_API_BASE_URL au build via DefinePlugin au lieu de dépendre seulement de process.env.",
            "Lire d'abord la valeur build-time dans api.ts, puis conserver la logique de validation existante.",
            "Bloquer npm run android / ios tant qu'un environnement explicite n'est pas choisi.",
            "Documenter l'usage android:local, android:shared-vps, build:*:prod.",
        ],
        "delivered": [
            "Ajout de __NS_API_BASE_URL__ et de l'injection webpack.",
            "Script require-explicit-env.js pour éviter les faux démarrages.",
            "Correction dédiée poussée sur GitHub: fix/frontend-nativescript-env-runtime.",
        ],
        "impact": [
            "Moins de faux timeouts côté mobile: l'application vise enfin la bonne URL au runtime.",
            "L'expérience développeur est plus fiable, surtout sur Android emulator et VPS partagé.",
            "Cette story complète la story 4 en supprimant un bug de runtime spécifique à NativeScript.",
        ],
        "files": [
            "frontend/webpack.config.js",
            "frontend/app/utils/api.ts",
            "frontend/package.json",
            "frontend/scripts/require-explicit-env.js",
            "branche: fix/frontend-nativescript-env-runtime",
        ],
        "evidence": [
            "npm run android échoue désormais immédiatement avec une consigne claire.",
            "La valeur NS_API_BASE_URL est injectée dans le bundle avant le lancement Android.",
        ],
    },
]

EXTRA_SLIDES = [
    {
        "title": "Travaux transverses souvent oubliés mais essentiels",
        "bullets": [
            "Documentation mise à jour sur plusieurs surfaces: README, DEPLOYMENT, PRISMA_WORKFLOW, ARCHITECTURE, fichiers .env.example.",
            "Uniformisation du langage d'exploitation: migrations Prisma, healthcheck, environnements frontend, variables de prod.",
            "Amélioration de la lisibilité globale du repo pour une équipe étudiante: scripts nommés, responsabilités séparées, messages d'erreur actionnables.",
            "Capitalisation d'incidents réels de déploiement pour affiner les garde-fous déjà en place.",
        ],
        "side_title": "Fichiers documentaires",
        "side_bullets": [
            "README.md",
            "DEPLOYMENT.md",
            "backend/PRISMA_WORKFLOW.md",
            "frontend/ARCHITECTURE.md",
            ".env.example / .env.prod.example / .env.local.example",
        ],
    },
    {
        "title": "Incident de recette VPS: preuve que les garde-fous servent",
        "bullets": [
            "Observation initiale: backend unhealthy -> nginx ne pouvait pas démarrer.",
            "Diagnostic structuré: logs backend, prisma:status, prisma:check, prisma:migrate:deploy.",
            "Conclusion: la base était saine, mais JWT_SECRET manquait dans backend/.env.prod.",
            "Résultat positif: le backend a refusé de démarrer proprement au lieu de tourner dans un état non sécurisé.",
        ],
        "side_title": "Leçons de sprint",
        "side_bullets": [
            "Les garde-fous de prod ont détecté la vraie cause.",
            "Le healthcheck et les logs structurés ont guidé le diagnostic.",
            "La documentation de déploiement doit inclure les variables réellement obligatoires.",
            "La revue de sprint peut montrer une amélioration concrète et vérifiable.",
        ],
    },
]


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    font_size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Aptos",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    paragraph.alignment = align
    return box


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    bullets,
    font_size=15,
    color=TEXT,
    bullet_color=PRIMARY,
    spacing=4,
    font_name="Aptos",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        paragraph.level = 0
        paragraph.space_after = Pt(spacing)
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.name = font_name
        if paragraph.runs:
            paragraph.runs[0].font.color.rgb = color
        paragraph._pPr.insert(
            0,
            paragraph._element.makeelement(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}buClr"
            ),
        )
    return box


def add_card_title(slide, x, y, title, accent=PRIMARY):
    add_textbox(slide, x, y, 2.9, 0.3, title, font_size=15, color=accent, bold=True)


def add_footer(slide, index):
    add_textbox(
        slide,
        0.55,
        7.05,
        12.2,
        0.22,
        f"Sprint 8 • Projet_Mobile_developpement_avancee • slide {index}",
        font_size=9,
        color=MUTED,
    )


def add_status_chip(slide, x, y, text):
    fg, bg = STATUS_STYLES.get(text, (INFO, INFO_BG))
    shape = add_rect(slide, x, y, 2.25, 0.42, bg, bg, radius=True)
    frame = shape.text_frame
    frame.text = ""
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Aptos"
    paragraph.alignment = PP_ALIGN.CENTER


def add_header_band(slide, title, subtitle):
    add_rect(slide, 0, 0, WIDE_W, 1.1, NAVY)
    add_rect(slide, 0, 0, 2.7, 1.1, PRIMARY)
    add_textbox(slide, 0.55, 0.2, 7.5, 0.34, title, font_size=26, color=WHITE, bold=True)
    add_textbox(slide, 0.55, 0.62, 7.8, 0.22, subtitle, font_size=11, color=WHITE)


def add_story_slide(prs: Presentation, story: dict, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, f"Story {story['number']} • {story['title']}", "Problème, sous-tâches, livrables, impact")
    add_status_chip(slide, 10.45, 0.28, story["status"])

    number_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(1.35), Inches(0.65), Inches(0.65))
    number_circle.fill.solid()
    number_circle.fill.fore_color.rgb = PRIMARY
    number_circle.line.color.rgb = PRIMARY
    number_circle.text_frame.text = story["number"]
    p = number_circle.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE

    add_rect(slide, 0.45, 1.25, 3.85, 2.25, WHITE, LIGHT, radius=True)
    add_rect(slide, 4.55, 1.25, 4.2, 2.25, WHITE, LIGHT, radius=True)
    add_rect(slide, 8.95, 1.25, 3.95, 2.25, WHITE, LIGHT, radius=True)
    add_rect(slide, 0.45, 3.75, 6.0, 2.65, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.65, 3.75, 6.25, 2.65, WHITE, LIGHT, radius=True)

    add_card_title(slide, 0.7, 1.45, "Problème traité")
    add_bullets(slide, 0.7, 1.8, 3.3, 1.45, story["problem"], font_size=14)

    add_card_title(slide, 4.8, 1.45, "Sous-tâches")
    add_bullets(slide, 4.8, 1.8, 3.7, 1.45, story["subtasks"], font_size=13)

    add_card_title(slide, 9.2, 1.45, "Valeur produite")
    add_bullets(slide, 9.2, 1.8, 3.45, 1.45, story["impact"], font_size=13)

    add_card_title(slide, 0.7, 3.95, "Livrables concrets")
    add_bullets(slide, 0.7, 4.3, 3.0, 1.6, story["delivered"], font_size=13)
    add_card_title(slide, 3.75, 3.95, "Fichiers clés")
    add_bullets(slide, 3.75, 4.3, 2.45, 1.6, story["files"], font_size=12, color=SLATE, bullet_color=SLATE)

    add_card_title(slide, 6.9, 3.95, "Preuves / signaux")
    add_bullets(slide, 6.9, 4.3, 5.7, 1.65, story["evidence"], font_size=12.5)

    add_footer(slide, slide_index)


def add_title_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, 0, 0, 4.5, 7.5, PRIMARY_DARK)
    add_rect(slide, 3.25, 0, 0.3, 7.5, PRIMARY)
    add_rect(slide, 4.2, 0.7, 8.4, 6.0, NAVY_LIGHT, NAVY_LIGHT, radius=True)

    add_textbox(slide, 0.65, 0.55, 2.6, 0.4, "Sprint 8 Review", font_size=16, color=WHITE, bold=True)
    add_textbox(slide, 0.65, 1.05, 2.2, 0.7, "Projet Mobile\ndéveloppement\navancé", font_size=24, color=WHITE, bold=True)
    add_textbox(
        slide,
        4.65,
        1.0,
        7.2,
        1.1,
        "Fiabiliser le backend, rendre le mobile explicite côté API et sécuriser la production",
        font_size=26,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        4.65,
        2.25,
        7.0,
        0.85,
        "Présentation structurée par stories, sous-tâches, livrables et impacts\navec distinction entre lots mergés et branches prêtes à merger.",
        font_size=14,
        color=RGBColor(226, 232, 240),
    )

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(9.7), Inches(0.9), height=Inches(1.05))

    add_rect(slide, 4.65, 3.25, 2.3, 1.1, WHITE, WHITE, radius=True)
    add_rect(slide, 7.15, 3.25, 2.3, 1.1, WHITE, WHITE, radius=True)
    add_rect(slide, 9.65, 3.25, 2.3, 1.1, WHITE, WHITE, radius=True)

    add_textbox(slide, 4.95, 3.48, 1.8, 0.22, "8 stories mergées", font_size=13, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 4.95, 3.76, 1.8, 0.25, "sur origin/main", font_size=11, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.45, 3.48, 1.8, 0.22, "2 branches prêtes", font_size=13, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.45, 3.76, 1.8, 0.25, "à merger", font_size=11, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.95, 3.48, 1.8, 0.22, "72 tests backend", font_size=13, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.95, 3.76, 1.8, 0.25, "validés localement", font_size=11, color=TEXT, align=PP_ALIGN.CENTER)

    add_rect(slide, 4.65, 4.75, 7.2, 1.2, RGBColor(15, 23, 42), RGBColor(15, 23, 42), radius=True)
    add_textbox(
        slide,
        4.95,
        5.0,
        6.6,
        0.65,
        "Axes majeurs: robustesse data Prisma, hardening prod, CI fiable, home testable,\nobservabilité exploitable, sécurité HTTP, runtime mobile explicite.",
        font_size=13,
        color=WHITE,
    )

    add_footer(slide, slide_index)


def add_agenda_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Vision du sprint", "Découpage par lots métier et techniques")

    cards = [
        ("Axe 1 • Fondations prod", ["Prisma", "Docker runtime", "env.ts", "variables critiques"]),
        ("Axe 2 • Qualité de livraison", ["CI Node 20", "Prisma en CI", "coverage stable", "tests ciblés"]),
        ("Axe 3 • Maintenabilité backend", ["home service", "repository", "tests unitaires", "structure modulaire"]),
        ("Axe 4 • Runtime mobile fiable", ["API explicite", "env par contexte", "NativeScript runtime", "VPS partagé"]),
    ]

    positions = [(0.65, 1.6), (6.85, 1.6), (0.65, 4.0), (6.85, 4.0)]
    for (title, bullets), (x, y) in zip(cards, positions):
        add_rect(slide, x, y, 5.8, 1.9, WHITE, LIGHT, radius=True)
        add_textbox(slide, x + 0.25, y + 0.2, 5.2, 0.25, title, font_size=17, color=PRIMARY, bold=True)
        add_bullets(slide, x + 0.25, y + 0.58, 5.2, 1.05, bullets, font_size=13)

    add_footer(slide, slide_index)


def add_summary_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Synthèse de sprint", "Ce que la revue doit retenir en 30 secondes")

    for idx, (title, subtitle) in enumerate(SUMMARY_CARDS):
        x = 0.65 + (idx % 2) * 6.1
        y = 1.55 + (idx // 2) * 1.6
        add_rect(slide, x, y, 5.4, 1.25, WHITE, LIGHT, radius=True)
        add_textbox(slide, x + 0.25, y + 0.23, 5.0, 0.25, title, font_size=22, color=PRIMARY, bold=True)
        add_textbox(slide, x + 0.25, y + 0.62, 4.9, 0.35, subtitle, font_size=12, color=TEXT)

    add_rect(slide, 0.65, 4.95, 12.0, 1.25, WHITE, LIGHT, radius=True)
    add_textbox(slide, 0.95, 5.15, 2.0, 0.22, "État d'intégration", font_size=16, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        0.95,
        5.45,
        11.2,
        0.55,
        [
            "Mergé sur main: stories 01 à 08.",
            "Branches dédiées prêtes à merger: story 09 (sécurité HTTP) et story 10 (runtime NativeScript / fiabilité API).",
            "Le sprint produit un socle plus sûr pour la prod et une meilleure vitesse de diagnostic en recette.",
        ],
        font_size=13,
    )

    add_footer(slide, slide_index)


def add_roadmap_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Roadmap des stories du sprint 8", "Lecture rapide du périmètre livré")

    y = 1.45
    for story in STORIES:
        fg, bg = STATUS_STYLES.get(story["status"], (INFO, INFO_BG))
        add_rect(slide, 0.7, y, 0.6, 0.34, PRIMARY, PRIMARY, radius=True)
        add_textbox(slide, 0.73, y + 0.04, 0.5, 0.2, story["number"], font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 1.45, y - 0.01, 6.9, 0.25, story["title"], font_size=16, color=TEXT, bold=True)
        chip = add_rect(slide, 9.7, y - 0.02, 2.3, 0.38, bg, bg, radius=True)
        chip_frame = chip.text_frame
        chip_frame.text = story["status"]
        chip_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        chip_frame.paragraphs[0].runs[0].font.size = Pt(11)
        chip_frame.paragraphs[0].runs[0].font.bold = True
        chip_frame.paragraphs[0].runs[0].font.color.rgb = fg
        y += 0.48

    add_footer(slide, slide_index)


def add_extra_slide(prs: Presentation, data: dict, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, data["title"], "Compléments importants pour la soutenance")
    add_rect(slide, 0.65, 1.55, 7.4, 4.8, WHITE, LIGHT, radius=True)
    add_rect(slide, 8.3, 1.55, 4.35, 4.8, WHITE, LIGHT, radius=True)
    add_card_title(slide, 0.9, 1.75, "Message clé")
    add_bullets(slide, 0.9, 2.1, 6.85, 3.8, data["bullets"], font_size=15)
    add_card_title(slide, 8.55, 1.75, data["side_title"])
    add_bullets(slide, 8.55, 2.1, 3.8, 3.8, data["side_bullets"], font_size=14)
    add_footer(slide, slide_index)


def add_closing_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, 0.7, 0.9, 11.95, 5.8, NAVY_LIGHT, NAVY_LIGHT, radius=True)
    add_textbox(slide, 1.1, 1.25, 6.6, 0.4, "Conclusion de sprint", font_size=30, color=WHITE, bold=True)
    add_textbox(
        slide,
        1.1,
        1.8,
        10.5,
        0.7,
        "Le sprint 8 a surtout transformé le projet en base plus déployable, plus sûre et plus testable.\nLa valeur n'est pas seulement dans le code livré, mais dans la réduction du risque opérationnel.",
        font_size=16,
        color=RGBColor(226, 232, 240),
    )

    add_rect(slide, 1.0, 3.0, 5.1, 2.5, WHITE, WHITE, radius=True)
    add_rect(slide, 6.35, 3.0, 5.6, 2.5, WHITE, WHITE, radius=True)
    add_textbox(slide, 1.25, 3.2, 4.5, 0.25, "Ce qui est démontrable maintenant", font_size=17, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        1.25,
        3.6,
        4.45,
        1.45,
        [
            "Backend qui refuse les mauvais déploiements.",
            "CI reproductible et lisible.",
            "Healthcheck utile pour Docker et recette.",
            "Home feed testable et parcours critiques couverts.",
        ],
        font_size=13,
    )
    add_textbox(slide, 6.6, 3.2, 5.0, 0.25, "Étapes suivantes recommandées", font_size=17, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        6.6,
        3.6,
        4.9,
        1.45,
        [
            "Merger la branche sécurité HTTP.",
            "Merger la branche runtime NativeScript env.",
            "Renseigner définitivement JWT_SECRET et CORS_ALLOWED_ORIGINS en prod.",
            "Pointer le mobile vers l'URL Nginx finale du VPS ou le futur domaine HTTPS.",
        ],
        font_size=13,
    )

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.65), Inches(0.95), height=Inches(0.95))

    add_footer(slide, slide_index)


def save_presentation(prs: Presentation) -> Path:
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)

    slide_index = 1
    add_title_slide(prs, slide_index)
    slide_index += 1
    add_agenda_slide(prs, slide_index)
    slide_index += 1
    add_summary_slide(prs, slide_index)
    slide_index += 1
    add_roadmap_slide(prs, slide_index)
    slide_index += 1

    for story in STORIES:
        add_story_slide(prs, story, slide_index)
        slide_index += 1

    for extra in EXTRA_SLIDES:
        add_extra_slide(prs, extra, slide_index)
        slide_index += 1

    add_closing_slide(prs, slide_index)

    return save_presentation(prs)


if __name__ == "__main__":
    output = build_presentation()
    print(output)
