from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from generate_sprint_8_presentation import (
    ROOT,
    LOGO_PATH,
    WIDE_H,
    WIDE_W,
    PRIMARY,
    PRIMARY_DARK,
    NAVY,
    NAVY_LIGHT,
    SLATE,
    TEXT,
    MUTED,
    LIGHT,
    LIGHTER,
    WHITE,
    SUCCESS,
    SUCCESS_BG,
    READY,
    READY_BG,
    INFO,
    INFO_BG,
    set_background,
    add_rect,
    add_textbox,
    add_bullets,
    add_card_title,
    add_footer,
    add_header_band,
)


OUTPUT_PATH = ROOT / "presentation_sprint_8_visual_stats.pptx"

STORY_STATUS = {
    "Mergées sur main": 8,
    "Prêtes à merger": 2,
}

SPRINT_FACTS = [
    ("Stories travaillées", 10),
    ("Stories mergées", 8),
    ("Branches prêtes", 2),
    ("Tests backend verts", 72),
    ("Migrations Prisma", 8),
    ("Environnements mobile", 3),
]

RISK_REDUCTION = {
    "Drift Prisma": (5, 1),
    "Config prod": (5, 1),
    "CI instable": (4, 2),
    "Diagnostic backend": (4, 2),
    "API mobile ambiguë": (4, 2),
}

STORY_COMPLEXITY = {
    "Drift DB": 4,
    "Docker runtime": 3,
    "env.ts prod": 5,
    "Frontend env": 6,
    "CI": 2,
    "Home refactor": 4,
    "Tests": 4,
    "Observabilité": 4,
    "HTTP security": 5,
    "Runtime NS": 5,
}

DELIVERY_GATES = {
    "Prisma checks": 3,
    "CI jobs": 3,
    "Health checks": 2,
    "Envs frontend": 3,
    "Auth endpoints protégés": 5,
}


def add_connector_line(slide, x1, y1, x2, y2, color=MUTED, width=2.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_arrow_box(slide, x, y, w, h, text, fill, font_color=WHITE, font_size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    shape.text_frame.text = text
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if p.runs:
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = font_color
        p.runs[0].font.name = "Aptos"
    return shape


def add_chart(slide, chart_type, x, y, w, h, categories, series, title=None):
    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for name, values in series:
        chart_data.add_series(name, values)

    graphic_frame = slide.shapes.add_chart(
        chart_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        chart_data,
    )
    chart = graphic_frame.chart
    chart.has_legend = True if len(series) > 1 or chart_type in {
        XL_CHART_TYPE.PIE,
        XL_CHART_TYPE.DOUGHNUT,
    } else False
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        title_frame = chart.chart_title.text_frame
        if title_frame.paragraphs and title_frame.paragraphs[0].runs:
            title_frame.paragraphs[0].runs[0].font.size = Pt(14)
            title_frame.paragraphs[0].runs[0].font.bold = True
            title_frame.paragraphs[0].runs[0].font.color.rgb = TEXT

    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = Pt(10)
    labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    labels.show_value = True
    labels.show_category_name = chart_type in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}

    if chart_type not in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}:
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(226, 232, 240)

    return chart


def add_metric_card(slide, x, y, title, value, tone="primary"):
    if tone == "primary":
        fg, bg = PRIMARY, WHITE
    elif tone == "success":
        fg, bg = SUCCESS, WHITE
    elif tone == "ready":
        fg, bg = READY, WHITE
    else:
        fg, bg = INFO, WHITE
    add_rect(slide, x, y, 2.8, 1.12, bg, LIGHT, radius=True)
    add_textbox(slide, x + 0.18, y + 0.18, 2.4, 0.22, title, font_size=11.5, color=SLATE, bold=True)
    add_textbox(slide, x + 0.18, y + 0.48, 2.4, 0.34, str(value), font_size=24, color=fg, bold=True)


def add_cover(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, 0, 0, WIDE_W, 1.0, PRIMARY_DARK)
    add_rect(slide, 0, 1.0, 3.35, 6.5, PRIMARY)
    add_rect(slide, 3.65, 0.95, 8.95, 5.95, NAVY_LIGHT, NAVY_LIGHT, radius=True)

    add_textbox(slide, 0.55, 0.24, 4.8, 0.34, "Sprint 8 • Présentation visuelle avec statistiques et schémas", font_size=22, color=WHITE, bold=True)
    add_textbox(slide, 0.7, 1.55, 2.0, 1.25, "Revue\nde sprint\nvisuelle", font_size=31, color=WHITE, bold=True)
    add_textbox(
        slide,
        4.0,
        1.35,
        7.9,
        0.95,
        "Projet_Mobile_developpement_avancee\nProduction • Qualité • Tests • Observabilité • Frontend mobile",
        font_size=25,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        4.0,
        2.55,
        7.9,
        0.7,
        "Objectif: montrer clairement ce qui a été livré,\nce qui est prêt, et quel risque technique a été réduit.",
        font_size=15,
        color=RGBColor(226, 232, 240),
    )

    add_rect(slide, 4.0, 4.15, 7.9, 1.45, WHITE, WHITE, radius=True)
    add_textbox(slide, 4.3, 4.4, 2.2, 0.22, "Chiffres clés", font_size=15, color=PRIMARY, bold=True)
    add_textbox(slide, 4.3, 4.8, 7.1, 0.42, "10 stories • 8 mergées • 2 branches prêtes • 72 tests backend • 8 migrations Prisma • 3 environnements mobile", font_size=13.2, color=TEXT)

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.85), Inches(0.14), height=Inches(0.72))

    add_footer(slide, slide_index)


def add_dashboard(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Dashboard du sprint", "Vue d'ensemble en indicateurs visuels")

    positions = [
        (0.65, 1.45), (3.6, 1.45), (6.55, 1.45),
        (9.5, 1.45), (2.15, 2.75), (5.1, 2.75),
    ]
    tones = ["primary", "success", "ready", "info", "primary", "primary"]
    for (title, value), (x, y), tone in zip(SPRINT_FACTS, positions, tones):
        add_metric_card(slide, x, y, title, value, tone=tone)

    add_rect(slide, 0.65, 4.25, 5.8, 2.35, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.75, 4.25, 5.9, 2.35, WHITE, LIGHT, radius=True)
    add_chart(
        slide,
        XL_CHART_TYPE.DOUGHNUT,
        0.95,
        4.6,
        2.7,
        1.65,
        STORY_STATUS.keys(),
        [("Stories", STORY_STATUS.values())],
        "Répartition des stories",
    )
    add_bullets(
        slide,
        3.65,
        4.85,
        2.3,
        1.25,
        [
            "80% du lot déjà mergé sur main.",
            "20% en branches dédiées prêtes à merger.",
            "Le sprint est majoritairement intégré.",
        ],
        font_size=12.2,
    )

    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        7.05,
        4.65,
        5.1,
        1.7,
        ["Pass", "Fail"],
        [("Backend tests", [72, 0])],
        "Résultat de la suite backend",
    )
    add_footer(slide, slide_index)


def add_stats_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Statistiques consolidées", "Volumes et garde-fous activés pendant le sprint")

    add_rect(slide, 0.65, 1.45, 6.0, 2.4, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.85, 1.45, 5.8, 2.4, WHITE, LIGHT, radius=True)
    add_rect(slide, 0.65, 4.1, 12.0, 2.25, WHITE, LIGHT, radius=True)

    add_chart(
        slide,
        XL_CHART_TYPE.BAR_CLUSTERED,
        0.95,
        1.78,
        5.35,
        1.75,
        DELIVERY_GATES.keys(),
        [("Total", DELIVERY_GATES.values())],
        "Garde-fous et surfaces couvertes",
    )

    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        7.15,
        1.78,
        5.1,
        1.75,
        STORY_COMPLEXITY.keys(),
        [("Fichiers principaux touchés", STORY_COMPLEXITY.values())],
        "Complexité relative par story",
    )

    add_textbox(slide, 0.95, 4.35, 4.0, 0.22, "Lecture des chiffres", font_size=16, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        0.95,
        4.72,
        10.8,
        1.15,
        [
            "Le sprint n'est pas centré sur une seule feature visible, mais sur la fiabilité transverse du système.",
            "La plus forte densité de changement se concentre sur l'environnement, le frontend API, la sécurité HTTP et le runtime mobile.",
            "Le backend dispose désormais de plusieurs lignes de défense successives: Prisma, env, healthcheck, CI, logs et tests critiques.",
        ],
        font_size=13.4,
    )
    add_textbox(slide, 8.9, 5.9, 3.0, 0.2, "Statistiques basées sur données réelles du repo\net sur les fichiers principaux touchés.", font_size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)
    add_footer(slide, slide_index)


def add_risk_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Réduction du risque technique", "Évaluation technique avant / après sur une échelle de 1 à 5")

    add_rect(slide, 0.65, 1.45, 8.25, 4.95, WHITE, LIGHT, radius=True)
    add_rect(slide, 9.15, 1.45, 3.5, 4.95, WHITE, LIGHT, radius=True)

    categories = list(RISK_REDUCTION.keys())
    before = [pair[0] for pair in RISK_REDUCTION.values()]
    after = [pair[1] for pair in RISK_REDUCTION.values()]
    add_chart(
        slide,
        XL_CHART_TYPE.BAR_CLUSTERED,
        0.95,
        1.8,
        7.6,
        4.1,
        categories,
        [("Avant", before), ("Après", after)],
        "Score de risque (5 = critique, 1 = maîtrisé)",
    )

    add_card_title(slide, 9.4, 1.75, "À retenir")
    add_bullets(
        slide,
        9.4,
        2.12,
        2.95,
        2.6,
        [
            "Le gain le plus net concerne Prisma et la configuration de prod.",
            "Le risque de mauvaise cible API côté mobile a été fortement réduit.",
            "La CI et l'observabilité n'éliminent pas tout risque, mais le rendent plus détectable.",
        ],
        font_size=13,
    )
    add_textbox(
        slide,
        9.4,
        5.25,
        2.9,
        0.55,
        "Note: ces scores sont une estimation d'impact technique pour la présentation, pas une mesure métier brute.",
        font_size=10.5,
        color=MUTED,
    )
    add_footer(slide, slide_index)


def add_architecture_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Schéma d'architecture du sprint", "Comment les changements s'articulent de bout en bout")

    add_arrow_box(slide, 0.75, 2.2, 2.0, 0.8, "Frontend mobile\nNativeScript", PRIMARY)
    add_arrow_box(slide, 3.1, 2.2, 1.95, 0.8, "api.ts\nconfig explicite", INFO)
    add_arrow_box(slide, 5.4, 2.2, 2.0, 0.8, "HTTP backend\nExpress + health + logs", RGBColor(37, 99, 235))
    add_arrow_box(slide, 7.75, 2.2, 1.8, 0.8, "Services métier\nhome/auth/profile", SUCCESS)
    add_arrow_box(slide, 9.9, 2.2, 1.45, 0.8, "Prisma\nchecks", READY)
    add_arrow_box(slide, 11.65, 2.2, 1.0, 0.8, "Postgres", PRIMARY_DARK)

    add_connector_line(slide, 2.75, 2.6, 3.1, 2.6)
    add_connector_line(slide, 5.05, 2.6, 5.4, 2.6)
    add_connector_line(slide, 7.4, 2.6, 7.75, 2.6)
    add_connector_line(slide, 9.55, 2.6, 9.9, 2.6)
    add_connector_line(slide, 11.35, 2.6, 11.65, 2.6)

    add_arrow_box(slide, 2.2, 4.35, 2.3, 0.72, "CI Node 20", NAVY)
    add_arrow_box(slide, 4.95, 4.35, 2.35, 0.72, "prisma validate + migrate + check", NAVY)
    add_arrow_box(slide, 7.75, 4.35, 2.05, 0.72, "tests + coverage", NAVY)
    add_arrow_box(slide, 10.25, 4.35, 2.05, 0.72, "déploiement VPS", NAVY)
    add_connector_line(slide, 4.5, 4.7, 4.95, 4.7)
    add_connector_line(slide, 7.3, 4.7, 7.75, 4.7)
    add_connector_line(slide, 9.8, 4.7, 10.25, 4.7)

    add_rect(slide, 0.85, 5.65, 11.8, 0.55, WHITE, WHITE, radius=True)
    add_textbox(
        slide,
        1.05,
        5.82,
        11.3,
        0.22,
        "Le sprint relie désormais explicitement le mobile, le backend HTTP, la couche métier, Prisma et la chaîne CI/déploiement.",
        font_size=12.5,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, slide_index)


def add_prisma_prod_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Prisma + production", "Le cœur des garde-fous de cohérence")

    add_rect(slide, 0.65, 1.45, 4.2, 4.95, WHITE, LIGHT, radius=True)
    add_rect(slide, 5.05, 1.45, 3.4, 4.95, WHITE, LIGHT, radius=True)
    add_rect(slide, 8.65, 1.45, 4.0, 4.95, WHITE, LIGHT, radius=True)

    add_card_title(slide, 0.9, 1.75, "Chaîne de démarrage backend")
    add_arrow_box(slide, 1.1, 2.1, 3.2, 0.56, "validateRuntimeEnv()", INFO, font_size=13)
    add_arrow_box(slide, 1.1, 2.95, 3.2, 0.56, "assertDatabaseIsReady()", READY, font_size=13)
    add_arrow_box(slide, 1.1, 3.8, 3.2, 0.56, "createHttpApp()", SUCCESS, font_size=13)
    add_arrow_box(slide, 1.1, 4.65, 3.2, 0.56, "app.listen()", PRIMARY, font_size=13)
    add_connector_line(slide, 2.7, 2.66, 2.7, 2.95)
    add_connector_line(slide, 2.7, 3.51, 2.7, 3.8)
    add_connector_line(slide, 2.7, 4.36, 2.7, 4.65)

    add_card_title(slide, 5.3, 1.75, "Contrôles")
    add_bullets(
        slide,
        5.3,
        2.1,
        2.8,
        3.8,
        [
            "_prisma_migrations présente",
            "migrations locales appliquées",
            "aucune migration en échec",
            "pas de migration inconnue côté DB",
            "tables/colonnes conformes au schéma",
            "DATABASE_URL et JWT_SECRET obligatoires en prod",
        ],
        font_size=12.8,
    )

    add_card_title(slide, 8.9, 1.75, "Bénéfices")
    add_bullets(
        slide,
        8.9,
        2.1,
        3.45,
        3.8,
        [
            "Évite les déploiements silencieusement incohérents.",
            "Donne un message d'erreur immédiatement exploitable.",
            "A permis d'isoler sur VPS une vraie panne config: JWT_SECRET absent.",
            "Fait converger dev, CI et prod autour du même workflow Prisma.",
        ],
        font_size=13,
    )
    add_footer(slide, slide_index)


def add_ci_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "CI GitHub Actions", "Pipeline simplifié, cohérent et compatible Node 20")

    add_rect(slide, 0.65, 1.45, 12.0, 2.1, WHITE, LIGHT, radius=True)
    steps = [
        ("Checkout", INFO),
        ("npm ci", INFO),
        ("prisma generate", READY),
        ("prisma validate", READY),
        ("migrate deploy", READY),
        ("prisma check", SUCCESS),
        ("tests / coverage / type-check", PRIMARY),
    ]
    x = 0.9
    for idx, (label, color) in enumerate(steps):
        add_arrow_box(slide, x, 2.05, 1.5 if idx < 2 else 1.7, 0.56, label, color, font_size=12)
        if idx < len(steps) - 1:
            add_connector_line(slide, x + (1.5 if idx < 2 else 1.7), 2.33, x + (1.5 if idx < 2 else 1.7) + 0.15, 2.33)
        x += 1.75

    add_rect(slide, 0.65, 3.9, 5.85, 2.45, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.8, 3.9, 5.85, 2.45, WHITE, LIGHT, radius=True)

    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        0.95,
        4.3,
        5.2,
        1.55,
        ["Backend tests", "Backend coverage", "Frontend type-check"],
        [("Jobs / checks", [1, 1, 1])],
        "Répartition lisible des contrôles",
    )
    add_bullets(
        slide,
        7.1,
        4.35,
        5.0,
        1.6,
        [
            "Node 20 aligné partout.",
            "CI fidèle au runtime Prisma réel.",
            "Moins de doublons entre jobs.",
            "Coverage backend stabilisé avec c8.",
        ],
        font_size=13.5,
    )
    add_footer(slide, slide_index)


def add_tests_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Tests critiques backend", "Couvrir les parcours à risque au lieu de viser un chiffre vide")

    add_rect(slide, 0.65, 1.45, 5.75, 4.95, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.7, 1.45, 5.95, 4.95, WHITE, LIGHT, radius=True)

    add_chart(
        slide,
        XL_CHART_TYPE.PIE,
        0.95,
        1.9,
        2.8,
        2.0,
        ["Pass", "Fail"],
        [("Suite backend", [72, 0])],
        "Résultat global",
    )
    add_textbox(slide, 3.9, 2.0, 2.0, 0.25, "Parcours couverts", font_size=15, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        3.9,
        2.38,
        2.1,
        2.4,
        [
            "auth",
            "home",
            "profile",
            "reservations",
            "password reset",
            "health",
            "guards / erreurs",
        ],
        font_size=12.6,
    )

    add_card_title(slide, 7.0, 1.75, "Exemples de scénarios utiles")
    add_bullets(
        slide,
        7.0,
        2.12,
        5.25,
        3.9,
        [
            "Login, refresh et mise à jour de profil via HTTP.",
            "Conflit email -> 409 proprement géré.",
            "health -> 200 avec x-request-id et checks app/db.",
            "home -> réponse valide même sans données ou avec token invalide.",
            "forgot-password -> pas d'énumération de comptes.",
            "reset-password -> code réutilisé refusé.",
            "notifications / tutorials / endpoints protégés -> 401/403 vérifiés.",
        ],
        font_size=13.2,
    )
    add_footer(slide, slide_index)


def add_home_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Refacto du module home", "Passage d'un flux couplé à un flux testable et injectable")

    add_rect(slide, 0.65, 1.45, 5.8, 4.95, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.75, 1.45, 5.9, 4.95, WHITE, LIGHT, radius=True)

    add_card_title(slide, 0.95, 1.75, "Avant / Après")
    add_arrow_box(slide, 1.05, 2.2, 1.45, 0.6, "HTTP", PRIMARY, font_size=13)
    add_arrow_box(slide, 2.75, 2.2, 1.45, 0.6, "Métier", PRIMARY, font_size=13)
    add_arrow_box(slide, 4.45, 2.2, 1.45, 0.6, "Prisma", PRIMARY, font_size=13)
    add_connector_line(slide, 2.5, 2.5, 2.75, 2.5)
    add_connector_line(slide, 4.2, 2.5, 4.45, 2.5)
    add_textbox(slide, 1.0, 3.15, 4.9, 0.28, "Après refacto", font_size=15, color=SUCCESS, bold=True)
    add_arrow_box(slide, 1.05, 3.65, 1.45, 0.6, "Controller", INFO, font_size=13)
    add_arrow_box(slide, 2.75, 3.65, 1.45, 0.6, "Service", SUCCESS, font_size=13)
    add_arrow_box(slide, 4.45, 3.65, 1.45, 0.6, "Repository", READY, font_size=13)
    add_connector_line(slide, 2.5, 3.95, 2.75, 3.95)
    add_connector_line(slide, 4.2, 3.95, 4.45, 3.95)
    add_textbox(slide, 1.0, 4.8, 4.8, 0.8, "Résultat: responsabilités visibles, dépendances injectables, logique métier testée sans Express ni Prisma réel.", font_size=12.5, color=TEXT)

    add_card_title(slide, 7.0, 1.75, "Tests spécifiques home")
    add_bullets(
        slide,
        7.0,
        2.15,
        5.2,
        3.7,
        [
            "Feed enrichi à partir des rendez-vous et rappels.",
            "Retour d'un feed par défaut si l'utilisateur n'a pas de données.",
            "Fallback propre et logging si une dépendance plante.",
            "Helpers métiers testés: prénom, format date, reminder message.",
        ],
        font_size=13.5,
    )
    add_footer(slide, slide_index)


def add_observability_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Observabilité + healthcheck", "Lire le backend sans stack de monitoring lourde")

    add_rect(slide, 0.65, 1.45, 12.0, 2.25, WHITE, LIGHT, radius=True)
    add_rect(slide, 0.65, 4.0, 5.85, 2.35, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.8, 4.0, 5.85, 2.35, WHITE, LIGHT, radius=True)

    add_textbox(slide, 0.95, 1.75, 2.3, 0.22, "Cycle d'une requête", font_size=16, color=PRIMARY, bold=True)
    add_arrow_box(slide, 1.0, 2.28, 1.7, 0.56, "req", INFO, font_size=13)
    add_arrow_box(slide, 3.05, 2.28, 2.15, 0.56, "requestLogger", SUCCESS, font_size=13)
    add_arrow_box(slide, 5.55, 2.28, 2.15, 0.56, "route/controller", READY, font_size=13)
    add_arrow_box(slide, 8.05, 2.28, 1.95, 0.56, "errorHandler", PRIMARY, font_size=13)
    add_arrow_box(slide, 10.35, 2.28, 1.7, 0.56, "res", PRIMARY_DARK, font_size=13)
    add_connector_line(slide, 2.7, 2.56, 3.05, 2.56)
    add_connector_line(slide, 5.2, 2.56, 5.55, 2.56)
    add_connector_line(slide, 7.7, 2.56, 8.05, 2.56)
    add_connector_line(slide, 10.0, 2.56, 10.35, 2.56)
    add_textbox(slide, 0.95, 3.05, 11.1, 0.26, "Champs tracés: requestId, method, url, statusCode, durationMs, responseBytes, userId, ip, user-agent", font_size=12.5, color=TEXT)

    add_card_title(slide, 0.95, 4.28, "Healthcheck retourné")
    add_bullets(
        slide,
        0.95,
        4.62,
        5.1,
        1.35,
        [
            "ok / status",
            "service",
            "environment",
            "timestamp",
            "checks.app.status",
            "checks.db.status",
        ],
        font_size=13.3,
    )

    add_card_title(slide, 7.1, 4.28, "Pourquoi c'est utile")
    add_bullets(
        slide,
        7.1,
        4.62,
        5.0,
        1.35,
        [
            "Réutilisable dans Docker, CI, curl et prod.",
            "Accélère le diagnostic des erreurs et lenteurs.",
            "Donne un état public utile sans exposer d'information sensible.",
        ],
        font_size=13.3,
    )
    add_footer(slide, slide_index)


def add_security_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Lot prêt: sécurité HTTP minimale", "Branche prête à merger")

    add_rect(slide, 0.65, 1.45, 5.95, 4.95, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.85, 1.45, 5.8, 4.95, WHITE, LIGHT, radius=True)

    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        0.95,
        1.9,
        5.1,
        1.8,
        ["Helmet", "CORS", "JSON limit", "Rate limit", "Public error"],
        [("Présent", [1, 1, 1, 1, 1])],
        "Garde-fous ajoutés",
    )
    add_bullets(
        slide,
        0.95,
        4.05,
        5.1,
        1.8,
        [
            "helmet compatible mobile",
            "allowlist CORS par environnement",
            "taille max JSON / urlencoded",
            "rate limit auth + reset",
            "messages d'erreur génériques masqués en prod",
        ],
        font_size=13,
    )

    add_card_title(slide, 7.1, 1.8, "Statut et intérêt")
    add_textbox(slide, 7.1, 2.15, 4.9, 0.35, "Cette story n'est pas encore mergée sur main,\nmais le lot est prêt sur la branche dédiée.", font_size=14, color=TEXT)
    add_bullets(
        slide,
        7.1,
        2.8,
        4.95,
        2.5,
        [
            "Réduit l'exposition aux abus sur login/register/reset.",
            "Améliore la posture prod sans casser le frontend mobile.",
            "Complète naturellement les stories env.ts + observabilité + CI.",
            "Tests dédiés déjà écrits pour sécuriser le merge.",
        ],
        font_size=13.3,
    )
    add_chip_box(slide, 7.1, 5.45, "Branche", "feat/backend-http-security-hardening")
    add_footer(slide, slide_index)


def add_chip_box(slide, x, y, label, value):
    add_textbox(slide, x, y, 1.0, 0.2, label, font_size=11, color=SLATE, bold=True)
    add_rect(slide, x, y + 0.26, 4.2, 0.45, INFO_BG, INFO_BG, radius=True)
    add_textbox(slide, x + 0.15, y + 0.35, 3.9, 0.16, value, font_size=10.5, color=INFO, bold=True)


def add_frontend_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Frontend mobile: environnements et runtime API", "Fin du fallback implicite et correction NativeScript runtime")

    add_rect(slide, 0.65, 1.45, 12.0, 2.15, WHITE, LIGHT, radius=True)
    add_rect(slide, 0.65, 3.9, 5.95, 2.45, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.85, 3.9, 5.8, 2.45, WHITE, LIGHT, radius=True)

    add_arrow_box(slide, 0.95, 2.1, 2.0, 0.58, ".env.local", INFO, font_size=13)
    add_arrow_box(slide, 3.25, 2.1, 2.2, 0.58, ".env.shared-vps", READY, font_size=13)
    add_arrow_box(slide, 5.75, 2.1, 1.9, 0.58, ".env.prod", PRIMARY, font_size=13)
    add_arrow_box(slide, 8.0, 2.1, 2.15, 0.58, "run-with-env.js", SUCCESS, font_size=13)
    add_arrow_box(slide, 10.5, 2.1, 1.8, 0.58, "api.ts", PRIMARY_DARK, font_size=13)
    add_connector_line(slide, 2.95, 2.39, 3.25, 2.39)
    add_connector_line(slide, 5.45, 2.39, 5.75, 2.39)
    add_connector_line(slide, 7.65, 2.39, 8.0, 2.39)
    add_connector_line(slide, 10.15, 2.39, 10.5, 2.39)
    add_textbox(slide, 0.95, 3.0, 11.1, 0.22, "Complément du lot: injection build-time de NS_API_BASE_URL dans le bundle NativeScript pour fiabiliser Android.", font_size=12.5, color=TEXT)

    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        0.95,
        4.3,
        5.1,
        1.55,
        ["Local", "VPS partagé", "Prod"],
        [("Contextes supportés", [1, 1, 1])],
        "Trois contextes explicitement supportés",
    )
    add_bullets(
        slide,
        7.1,
        4.25,
        4.95,
        1.8,
        [
            "Plus de fallback IP caché vers le VPS.",
            "npm run android / ios sont bloqués sans contexte explicite.",
            "Correction runtime NativeScript poussée sur branche dédiée.",
            "L'équipe choisit consciemment sa cible API.",
        ],
        font_size=13.3,
    )
    add_footer(slide, slide_index)


def add_incident_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Cas concret: incident VPS", "Le sprint a été utile en situation réelle de recette")

    add_rect(slide, 0.65, 1.45, 12.0, 4.95, WHITE, LIGHT, radius=True)
    steps = [
        ("1", "backend unhealthy", READY_BG, READY),
        ("2", "logs backend", INFO_BG, INFO),
        ("3", "Prisma OK", SUCCESS_BG, SUCCESS),
        ("4", "JWT_SECRET missing", RGBColor(254, 226, 226), PRIMARY),
        ("5", "cause isolée", RGBColor(224, 231, 255), RGBColor(79, 70, 229)),
    ]
    x = 0.95
    for idx, (num, label, bg, fg) in enumerate(steps):
        add_rect(slide, x, 2.35, 1.95, 1.25, bg, bg, radius=True)
        add_textbox(slide, x + 0.15, 2.56, 1.65, 0.22, num, font_size=20, color=fg, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 2.95, 1.75, 0.35, label, font_size=12.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_connector_line(slide, x + 1.95, 2.98, x + 2.2, 2.98)
        x += 2.3

    add_bullets(
        slide,
        1.0,
        4.45,
        11.0,
        1.2,
        [
            "Le preflight Prisma a confirmé que la base était saine.",
            "validateRuntimeEnv() a immédiatement signalé la vraie panne: JWT_SECRET absent.",
            "Conclusion: les garde-fous ajoutés ne sont pas théoriques, ils ont servi pendant la recette VPS.",
        ],
        font_size=14,
    )
    add_footer(slide, slide_index)


def add_roadmap_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, 0.75, 0.95, 11.85, 5.85, NAVY_LIGHT, NAVY_LIGHT, radius=True)
    add_textbox(slide, 1.15, 1.25, 5.0, 0.38, "Conclusion visuelle", font_size=29, color=WHITE, bold=True)
    add_textbox(
        slide,
        1.15,
        1.8,
        10.4,
        0.7,
        "Le sprint 8 a surtout professionnalisé la plateforme.\nLa démonstration ne repose pas seulement sur une feature, mais sur une chaîne complète plus fiable.",
        font_size=16,
        color=RGBColor(226, 232, 240),
    )

    add_rect(slide, 1.1, 3.0, 5.25, 2.55, WHITE, WHITE, radius=True)
    add_rect(slide, 6.65, 3.0, 5.05, 2.55, WHITE, WHITE, radius=True)
    add_textbox(slide, 1.38, 3.22, 4.4, 0.25, "Ce qu'on peut montrer au jury", font_size=16, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        1.38,
        3.58,
        4.45,
        1.5,
        [
            "des statistiques claires",
            "une architecture plus propre",
            "des tests utiles et verts",
            "des incidents de prod mieux diagnostiqués",
            "un mobile qui cible explicitement son API",
        ],
        font_size=13.2,
    )
    add_textbox(slide, 6.92, 3.22, 4.1, 0.25, "Actions suivantes", font_size=16, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        6.92,
        3.58,
        4.2,
        1.5,
        [
            "merger les lots story 09 et 10",
            "finaliser backend/.env.prod sur VPS",
            "ouvrir l'URL Nginx finale au frontend mobile",
            "garder ce niveau de garde-fous sur les prochains sprints",
        ],
        font_size=13.2,
    )

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.95), Inches(1.08), height=Inches(0.82))
    add_footer(slide, slide_index)


def save_presentation(prs: Presentation) -> Path:
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)

    slide_index = 1
    add_cover(prs, slide_index)
    slide_index += 1
    add_dashboard(prs, slide_index)
    slide_index += 1
    add_stats_slide(prs, slide_index)
    slide_index += 1
    add_risk_slide(prs, slide_index)
    slide_index += 1
    add_architecture_slide(prs, slide_index)
    slide_index += 1
    add_prisma_prod_slide(prs, slide_index)
    slide_index += 1
    add_ci_slide(prs, slide_index)
    slide_index += 1
    add_tests_slide(prs, slide_index)
    slide_index += 1
    add_home_slide(prs, slide_index)
    slide_index += 1
    add_observability_slide(prs, slide_index)
    slide_index += 1
    add_security_slide(prs, slide_index)
    slide_index += 1
    add_frontend_slide(prs, slide_index)
    slide_index += 1
    add_incident_slide(prs, slide_index)
    slide_index += 1
    add_roadmap_slide(prs, slide_index)

    return save_presentation(prs)


if __name__ == "__main__":
    output = build_presentation()
    print(output)
