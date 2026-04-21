from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from generate_sprint_8_presentation import (
    ROOT,
    LOGO_PATH,
    WIDE_H,
    WIDE_W,
    PRIMARY,
    PRIMARY_DARK,
    NAVY,
    NAVY_LIGHT,
    SLATE,
    TEXT,
    MUTED,
    LIGHT,
    LIGHTER,
    WHITE,
    SUCCESS,
    READY,
    INFO,
    SUCCESS_BG,
    READY_BG,
    INFO_BG,
    STATUS_STYLES,
    set_background,
    add_rect,
    add_textbox,
    add_bullets,
    add_card_title,
    add_footer,
    add_header_band,
    add_status_chip,
)


OUTPUT_PATH = ROOT / "presentation_sprint_8_review_technique.pptx"

TECH_STORIES = [
    {
        "number": "01",
        "title": "Corriger le drift Prisma / base de données",
        "status": "Mergée sur main",
        "goal": [
            "Empêcher tout démarrage backend sur une base qui ne correspond pas au code déployé.",
            "Transformer les incohérences Prisma en erreur de préflight claire et actionnable.",
        ],
        "subtasks": [
            "Scanner les migrations locales et l'état réel de _prisma_migrations.",
            "Comparer tables/colonnes réelles avec le schéma Prisma attendu.",
            "Bloquer le boot si migration absente, échouée ou inconnue.",
        ],
        "implementation": [
            "Ajout de databaseReadiness.ts + schemaManifest.ts.",
            "Préflight lancé dans server.ts avant createHttpApp().",
            "Commande prisma:check réutilisable en CI, local et prod.",
        ],
        "files": [
            "backend/src/data/prisma/databaseReadiness.ts",
            "backend/src/data/prisma/schemaManifest.ts",
            "backend/src/server.ts",
            "backend/PRISMA_WORKFLOW.md",
        ],
        "validation": [
            "Erreurs explicites observées sur table _prisma_migrations absente.",
            "VPS validé avec prisma:status, prisma:check, prisma:migrate:deploy.",
        ],
    },
    {
        "number": "02",
        "title": "Retirer les secrets du Dockerfile",
        "status": "Mergée sur main",
        "goal": [
            "Séparer proprement build-time et runtime dans l'image backend.",
            "Éviter tout secret ou configuration de prod embarqué dans l'image.",
        ],
        "subtasks": [
            "Conserver uniquement une DATABASE_URL placeholder pour le build Prisma.",
            "Injecter la vraie configuration via Docker Compose au runtime.",
            "Préserver prisma migrate deploy au démarrage du conteneur.",
        ],
        "implementation": [
            "Dockerfile multi-stage simplifié.",
            "CMD runtime: vérification DATABASE_URL puis migrate deploy puis node dist/server.js.",
            "Documentation de déploiement mise à jour.",
        ],
        "files": [
            "backend/Dockerfile",
            "docker-compose.prod.yml",
            "DEPLOYMENT.md",
        ],
        "validation": [
            "Le build ne dépend plus d'un secret réel.",
            "Le runtime échoue explicitement si DATABASE_URL manque.",
        ],
    },
    {
        "number": "03",
        "title": "Durcir env.ts pour la production",
        "status": "Mergée sur main",
        "goal": [
            "Supprimer les fallbacks dangereux et rendre la prod explicitement validée.",
            "Centraliser les règles d'environnement dans une seule couche lisible.",
        ],
        "subtasks": [
            "Ajouter requireEnv() et requireEnvInProduction().",
            "Exiger DATABASE_URL et JWT_SECRET en production.",
            "Désactiver DEMO_MODE par défaut.",
            "Valider aussi les origines CORS et la config associée.",
        ],
        "implementation": [
            "validateRuntimeEnv() appelé au démarrage.",
            "Messages d'erreur actionnables et bloquants.",
            "backend/.env.example et .env.prod.example réalignés.",
        ],
        "files": [
            "backend/src/config/env.ts",
            "backend/.env.example",
            "backend/.env.prod.example",
            "README.md",
            "DEPLOYMENT.md",
        ],
        "validation": [
            "Cas réel observé: JWT_SECRET manquant sur VPS, backend stoppé proprement.",
            "La cause a été trouvée sans ambiguïté grâce au préflight runtime.",
        ],
    },
    {
        "number": "04",
        "title": "Remplacer le fallback d'IP VPS dans le frontend mobile",
        "status": "Mergée sur main",
        "goal": [
            "Éviter qu'un développeur pointe vers la prod sans le vouloir.",
            "Rendre explicites les contextes local, VPS partagé et prod.",
        ],
        "subtasks": [
            "Supprimer le fallback implicite d'IP dans api.ts.",
            "Normaliser l'URL API et rejeter les mauvaises valeurs.",
            "Introduire des fichiers .env dédiés et des scripts de run explicites.",
        ],
        "implementation": [
            "Lecture stricte de NS_API_BASE_URL.",
            "Scripts android:local / android:shared-vps / build:*:prod.",
            "Documentation frontend mise à jour.",
        ],
        "files": [
            "frontend/app/utils/api.ts",
            "frontend/package.json",
            "frontend/scripts/run-with-env.js",
            "frontend/.env.local.example",
            "frontend/.env.shared-vps.example",
            "frontend/.env.prod.example",
        ],
        "validation": [
            "L'app n'a plus de fallback VPS hardcodé.",
            "Le message [API_CONFIG] guide directement vers la bonne commande.",
        ],
    },
    {
        "number": "05",
        "title": "Nettoyer et fiabiliser la CI GitHub Actions",
        "status": "Mergée sur main",
        "goal": [
            "Rendre la CI stable, lisible et alignée avec le vrai pipeline backend/frontend.",
            "Vérifier Prisma avant les tests au lieu de découvrir les problèmes plus tard.",
        ],
        "subtasks": [
            "Aligner Node sur 20 pour tous les jobs.",
            "Valider Prisma et appliquer les migrations sur une base CI propre.",
            "Exécuter prisma:check avant build/tests.",
            "Stabiliser la couverture backend avec c8.",
        ],
        "implementation": [
            "Workflow CI avec matrix backend tests/coverage.",
            "Frontend limité au type-check ciblé et rapide.",
            "Suppression du schéma coverage incompatible Node 20.",
        ],
        "files": [
            ".github/workflows/ci.yml",
            "backend/package.json",
        ],
        "validation": [
            "Ordre final: install -> prisma generate -> validate -> migrate -> check -> test.",
            "Script test:coverage:ci réparé pour la CI Node 20.",
        ],
    },
    {
        "number": "06",
        "title": "Refactoriser le module home pour le rendre testable",
        "status": "Mergée sur main",
        "goal": [
            "Déplacer la logique métier hors du contrôleur HTTP.",
            "Pouvoir tester le home feed sans dépendre de Prisma ou d'Express.",
        ],
        "subtasks": [
            "Créer home.service avec logique pure et injection de dépendances.",
            "Déplacer l'accès aux données dans home.repository.",
            "Conserver une compatibilité API maximale.",
        ],
        "implementation": [
            "Helpers métiers: displayName, appointment label, reminder message.",
            "Fallback propre si repository indisponible.",
            "home.controller réduit à la couche transport.",
        ],
        "files": [
            "backend/src/modules/home/home.controller.ts",
            "backend/src/modules/home/home.service.ts",
            "backend/src/modules/home/home.repository.ts",
            "backend/tests/homeService.test.ts",
        ],
        "validation": [
            "Cas testés: enrichi, vide, dépendance en erreur.",
            "Le module est devenu lisible et injectable.",
        ],
    },
    {
        "number": "07",
        "title": "Ajouter des tests vraiment utiles sur les routes critiques",
        "status": "Mergée sur main",
        "goal": [
            "Couvrir les scénarios métier et sécurité, pas seulement un taux abstrait de coverage.",
            "Sécuriser auth, home, reservations, profile et password reset.",
        ],
        "subtasks": [
            "Créer une base de test PostgreSQL temporaire et appliquer les migrations à chaud.",
            "Tester les parcours succès, erreurs métier et accès interdits.",
            "Vérifier les réponses publiques, les headers utiles et les routes protégées.",
        ],
        "implementation": [
            "httpIntegration.test.ts monte l'app sur un port éphémère.",
            "applyMigrations() lancé avant les scénarios d'intégration.",
            "Tests additionnels sur health, home service, guards et helpers.",
        ],
        "files": [
            "backend/tests/httpIntegration.test.ts",
            "backend/tests/homeService.test.ts",
            "backend/tests/health.test.ts",
            "backend/tests/controllerInternals.test.ts",
        ],
        "validation": [
            "Suite locale exécutée: 72 tests verts.",
            "Exemples couverts: 401, 403, 404, 409, home vide valide, reset code réutilisé rejeté.",
        ],
    },
    {
        "number": "08",
        "title": "Ajouter observabilité + healthcheck utile",
        "status": "Mergée sur main",
        "goal": [
            "Pouvoir diagnostiquer rapidement les problèmes backend en dev, CI et prod.",
            "Rendre le healthcheck exploitable par Docker et un humain.",
        ],
        "subtasks": [
            "Ajouter request id, durée, taille de réponse, IP, userId et méthode dans les logs.",
            "Créer un /health avec statut app, statut DB, timestamp et environnement.",
            "Décrire les erreurs Prisma dans une structure lisible.",
        ],
        "implementation": [
            "Middleware requestLogger avec x-request-id.",
            "readHealthCheck() basé sur une sonde DB.",
            "describePrismaError() pour les logs structurés côté backend.",
        ],
        "files": [
            "backend/src/core/http/middleware/requestLogger.ts",
            "backend/src/core/http/health.ts",
            "backend/src/data/prisma/prismaError.ts",
            "backend/tests/health.test.ts",
        ],
        "validation": [
            "Health endpoint vérifié dans les tests d'intégration.",
            "Logs structurés visibles pendant npm test et sur le VPS.",
        ],
    },
    {
        "number": "09",
        "title": "Durcir la sécurité HTTP minimale du backend",
        "status": "Branche prête à merger",
        "goal": [
            "Ajouter des garde-fous HTTP simples mais utiles sans casser le mobile.",
            "Fermer les trous les plus visibles côté headers, CORS, taille de payload et abuse auth.",
        ],
        "subtasks": [
            "Ajouter helmet avec options compatibles NativeScript/mobile.",
            "Rendre CORS explicite en production via allowlist.",
            "Ajouter rate limit sur register/login/forgot/reset.",
            "Masquer les erreurs génériques en prod tout en gardant les AppError métier.",
        ],
        "implementation": [
            "Nouveau security.ts avec createCorsOptions() et createAuthRateLimit().",
            "Nouveau httpErrors.ts pour les messages publics et statuts HTTP.",
            "createHttpApp.ts durci avec trust proxy, json limit et helmet.",
        ],
        "files": [
            "backend/src/core/http/security.ts",
            "backend/src/core/http/httpErrors.ts",
            "backend/src/core/http/createHttpApp.ts",
            "backend/tests/httpSecurity.test.ts",
            "branche: feat/backend-http-security-hardening",
        ],
        "validation": [
            "Tests dédiés: CORS, status resolution, public error message.",
            "Le lot est prêt à merger mais pas encore sur main.",
        ],
    },
    {
        "number": "10",
        "title": "Optimiser la fiabilité perçue côté frontend API",
        "status": "Branche prête à merger",
        "goal": [
            "Supprimer les faux timeouts et mauvais ciblages d'API dus au runtime NativeScript.",
            "Rendre l'environnement frontend réellement effectif jusque dans le bundle Android.",
        ],
        "subtasks": [
            "Injecter NS_API_BASE_URL au build via DefinePlugin.",
            "Lire la valeur build-time dans api.ts avant process.env.",
            "Bloquer les scripts android/ios génériques tant qu'un environnement explicite n'est pas choisi.",
            "Documenter les usages local / VPS / prod pour l'équipe.",
        ],
        "implementation": [
            "webpack.config.js injecte __NS_API_BASE_URL__.",
            "require-explicit-env.js évite les faux démarrages.",
            "Fix poussé sur la branche fix/frontend-nativescript-env-runtime.",
        ],
        "files": [
            "frontend/webpack.config.js",
            "frontend/app/utils/api.ts",
            "frontend/package.json",
            "frontend/scripts/require-explicit-env.js",
            "branche: fix/frontend-nativescript-env-runtime",
        ],
        "validation": [
            "type-check frontend vert.",
            "La variable NS_API_BASE_URL est confirmée dans le bundle avant le run Android.",
        ],
    },
]


def add_chip(slide, x, y, text, fg, bg):
    shape = add_rect(slide, x, y, 1.65, 0.34, bg, bg, radius=True)
    frame = shape.text_frame
    frame.text = text
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    frame.paragraphs[0].runs[0].font.size = Pt(10)
    frame.paragraphs[0].runs[0].font.bold = True
    frame.paragraphs[0].runs[0].font.color.rgb = fg


def add_cover(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, 0, 0, 13.333, 1.15, PRIMARY_DARK)
    add_rect(slide, 0, 1.15, 3.1, 6.35, PRIMARY)
    add_rect(slide, 3.45, 0.95, 9.25, 5.95, NAVY_LIGHT, NAVY_LIGHT, radius=True)
    add_textbox(slide, 0.5, 0.26, 4.0, 0.35, "Sprint 8 • Version Enseignant / Technique", font_size=22, color=WHITE, bold=True)
    add_textbox(slide, 0.58, 1.55, 2.0, 1.45, "Revue\ntechnique\nstructurée", font_size=29, color=WHITE, bold=True)
    add_textbox(
        slide,
        3.9,
        1.45,
        7.9,
        1.05,
        "Projet_Mobile_developpement_avancee\nArchitecture, qualité, prod, CI, tests et exploitation",
        font_size=24,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        3.9,
        2.75,
        7.9,
        0.75,
        "Objectif du support: expliquer précisément ce qui a été changé,\npourquoi cela a été fait et comment cela réduit le risque technique.",
        font_size=15,
        color=RGBColor(226, 232, 240),
    )
    add_rect(slide, 3.95, 4.15, 2.45, 1.15, WHITE, WHITE, radius=True)
    add_rect(slide, 6.65, 4.15, 2.45, 1.15, WHITE, WHITE, radius=True)
    add_rect(slide, 9.35, 4.15, 2.45, 1.15, WHITE, WHITE, radius=True)
    add_textbox(slide, 4.2, 4.35, 2.0, 0.25, "Back + Infra", font_size=15, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 4.2, 4.7, 2.0, 0.22, "Prisma • Docker • env • health", font_size=10.5, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 6.9, 4.35, 2.0, 0.25, "Qualité", font_size=15, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 6.9, 4.7, 2.0, 0.22, "CI • coverage • tests critiques", font_size=10.5, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.6, 4.35, 2.0, 0.25, "Frontend mobile", font_size=15, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.6, 4.7, 2.0, 0.22, "env explicite • runtime NativeScript", font_size=10.5, color=TEXT, align=PP_ALIGN.CENTER)

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.9), Inches(0.18), height=Inches(0.7))

    add_footer(slide, slide_index)


def add_scope_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Périmètre du sprint 8", "Backlog structuré en stories et en sous-tâches techniques")
    add_rect(slide, 0.65, 1.45, 12.0, 1.0, WHITE, LIGHT, radius=True)
    add_textbox(
        slide,
        0.95,
        1.72,
        11.4,
        0.45,
        "Le sprint a principalement renforcé la robustesse du système plutôt que d'ajouter une seule fonctionnalité visible. "
        "Le thème central est la réduction du risque: drift Prisma, env prod, CI, sécurité HTTP, tests, observabilité et ciblage mobile de l'API.",
        font_size=14,
        color=TEXT,
    )

    groups = [
        ("Fondations prod", ["01 Drift Prisma", "02 Secrets Dockerfile", "03 env.ts prod"]),
        ("Qualité de livraison", ["05 CI GitHub", "07 Tests critiques", "08 Observabilité"]),
        ("Maintenabilité métier", ["06 Refacto home testable"]),
        ("Frontend API mobile", ["04 Fin du fallback VPS", "10 Runtime NativeScript"]),
        ("Sécurité HTTP", ["09 Hardening HTTP minimal"]),
    ]

    positions = [
        (0.65, 2.8, 3.7),
        (4.55, 2.8, 4.0),
        (8.75, 2.8, 3.9),
        (0.65, 5.0, 5.8),
        (6.75, 5.0, 5.9),
    ]
    for (title, bullets), (x, y, w) in zip(groups, positions):
        add_rect(slide, x, y, w, 1.55, WHITE, LIGHT, radius=True)
        add_textbox(slide, x + 0.22, y + 0.2, w - 0.44, 0.25, title, font_size=16, color=PRIMARY, bold=True)
        add_bullets(slide, x + 0.22, y + 0.55, w - 0.44, 0.75, bullets, font_size=12.5)

    add_footer(slide, slide_index)


def add_architecture_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, "Vue architecture des changements", "Lecture par couches du système")

    add_rect(slide, 0.7, 1.65, 12.0, 4.7, WHITE, LIGHT, radius=True)

    layers = [
        ("Frontend mobile", ["NativeScript/Vue", "api.ts centralisé", "environnements explicites", "runtime build-time env"]),
        ("HTTP backend", ["Express", "request logger", "error handler", "healthcheck", "security guards"]),
        ("Métier backend", ["home.service", "auth/profile/reservations", "règles métier testées"]),
        ("Data / Prisma", ["migrations versionnées", "databaseReadiness", "prisma:check", "prisma migrate deploy"]),
        ("Livraison", ["Docker runtime propre", "CI Node 20", "tests + coverage", "VPS/Nginx"]),
    ]

    y = 1.95
    colors = [INFO_BG, READY_BG, SUCCESS_BG, RGBColor(237, 233, 254), RGBColor(254, 249, 195)]
    accents = [INFO, READY, SUCCESS, RGBColor(109, 40, 217), RGBColor(161, 98, 7)]
    for idx, (title, bullets) in enumerate(layers):
        add_rect(slide, 1.0, y, 11.35, 0.72, colors[idx], colors[idx], radius=True)
        add_textbox(slide, 1.2, y + 0.11, 2.45, 0.22, title, font_size=16, color=accents[idx], bold=True)
        add_textbox(slide, 3.25, y + 0.09, 8.6, 0.32, " • ".join(bullets), font_size=11.5, color=TEXT)
        if idx < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.35), Inches(y + 0.73), Inches(0.45), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.color.rgb = MUTED
        y += 0.9

    add_textbox(
        slide,
        1.0,
        6.55,
        11.4,
        0.3,
        "Message clé: la majorité des changements sont transverses. Le sprint a consolidé la chaîne complète, du mobile jusqu'au déploiement.",
        font_size=12.5,
        color=SLATE,
    )
    add_footer(slide, slide_index)


def add_story_slide(prs: Presentation, story: dict, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHT)
    add_header_band(slide, f"Story {story['number']} • {story['title']}", "Version technique détaillée")
    add_status_chip(slide, 10.45, 0.28, story["status"])

    add_rect(slide, 0.55, 1.35, 3.7, 2.15, WHITE, LIGHT, radius=True)
    add_rect(slide, 4.45, 1.35, 4.0, 2.15, WHITE, LIGHT, radius=True)
    add_rect(slide, 8.65, 1.35, 4.1, 2.15, WHITE, LIGHT, radius=True)
    add_rect(slide, 0.55, 3.75, 4.4, 2.25, WHITE, LIGHT, radius=True)
    add_rect(slide, 5.15, 3.75, 7.6, 2.25, WHITE, LIGHT, radius=True)

    add_card_title(slide, 0.8, 1.56, "Objectif technique")
    add_bullets(slide, 0.8, 1.9, 3.15, 1.35, story["goal"], font_size=13.5)

    add_card_title(slide, 4.7, 1.56, "Sous-tâches")
    add_bullets(slide, 4.7, 1.9, 3.45, 1.35, story["subtasks"], font_size=13)

    add_card_title(slide, 8.9, 1.56, "Implémentation")
    add_bullets(slide, 8.9, 1.9, 3.55, 1.35, story["implementation"], font_size=13)

    add_card_title(slide, 0.8, 3.95, "Fichiers / modules touchés")
    add_bullets(slide, 0.8, 4.3, 3.85, 1.3, story["files"], font_size=12.3, color=SLATE, bullet_color=SLATE)

    add_card_title(slide, 5.4, 3.95, "Validation / preuves")
    add_bullets(slide, 5.4, 4.3, 6.95, 1.35, story["validation"], font_size=12.8)

    add_footer(slide, slide_index)


def add_incident_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, LIGHTER)
    add_header_band(slide, "Étude de cas recette VPS", "Comment les nouveaux garde-fous ont aidé en situation réelle")

    add_rect(slide, 0.7, 1.55, 5.95, 4.95, WHITE, LIGHT, radius=True)
    add_rect(slide, 6.85, 1.55, 5.75, 4.95, WHITE, LIGHT, radius=True)

    add_card_title(slide, 0.95, 1.8, "Chronologie du diagnostic")
    add_bullets(
        slide,
        0.95,
        2.15,
        5.4,
        3.8,
        [
            "1. Nginx ne démarre pas car backend unhealthy.",
            "2. Logs backend: migrations OK, mais validateRuntimeEnv() bloque.",
            "3. Diagnostic ciblé dans le conteneur: DATABASE_URL=present, JWT_SECRET=missing.",
            "4. Cause réelle isolée: variable obligatoire absente en production.",
        ],
        font_size=14,
    )

    add_card_title(slide, 7.1, 1.8, "Ce que cela prouve")
    add_bullets(
        slide,
        7.1,
        2.15,
        5.2,
        3.8,
        [
            "Le préflight Prisma n'a pas masqué la réalité: il a d'abord validé la couche data.",
            "Le hardening env.ts a empêché un backend mal sécurisé de démarrer.",
            "L'observabilité a permis d'écarter immédiatement une fausse piste réseau.",
            "Le support de sprint peut montrer non seulement du code, mais un bénéfice opérationnel réel.",
        ],
        font_size=14,
    )

    add_footer(slide, slide_index)


def add_closing_slide(prs: Presentation, slide_index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_rect(slide, 0.85, 0.95, 11.7, 5.85, NAVY_LIGHT, NAVY_LIGHT, radius=True)
    add_textbox(slide, 1.2, 1.25, 6.5, 0.4, "Conclusion technique", font_size=28, color=WHITE, bold=True)
    add_textbox(
        slide,
        1.2,
        1.8,
        10.2,
        0.7,
        "Le sprint 8 a surtout sécurisé la chaîne de livraison et d'exécution.\nLa valeur principale est la réduction des états incohérents, des déploiements fragiles et des diagnostics lents.",
        font_size=16,
        color=RGBColor(226, 232, 240),
    )

    add_rect(slide, 1.15, 3.0, 5.2, 2.55, WHITE, WHITE, radius=True)
    add_rect(slide, 6.65, 3.0, 5.0, 2.55, WHITE, WHITE, radius=True)
    add_textbox(slide, 1.42, 3.2, 4.5, 0.25, "Ce qui est désormais robuste", font_size=16, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        1.42,
        3.55,
        4.45,
        1.5,
        [
            "Base Prisma et schéma versionnés et vérifiés.",
            "Boot prod contrôlé par des variables critiques obligatoires.",
            "CI reproductible et centrée sur Prisma + tests.",
            "Frontend mobile ciblant explicitement son API.",
        ],
        font_size=13,
    )
    add_textbox(slide, 6.92, 3.2, 4.3, 0.25, "Prochain lot recommandé", font_size=16, color=PRIMARY, bold=True)
    add_bullets(
        slide,
        6.92,
        3.55,
        4.2,
        1.5,
        [
            "Merger les branches story 09 et story 10.",
            "Finaliser backend/.env.prod sur le VPS.",
            "Activer l'URL Nginx finale côté mobile partagé.",
            "Conserver cette discipline de préflight sur les prochains sprints.",
        ],
        font_size=13,
    )

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.85), Inches(1.1), height=Inches(0.85))

    add_footer(slide, slide_index)


def save_presentation(prs: Presentation) -> Path:
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)

    slide_index = 1
    add_cover(prs, slide_index)
    slide_index += 1
    add_scope_slide(prs, slide_index)
    slide_index += 1
    add_architecture_slide(prs, slide_index)
    slide_index += 1

    for story in TECH_STORIES:
        add_story_slide(prs, story, slide_index)
        slide_index += 1

    add_incident_slide(prs, slide_index)
    slide_index += 1
    add_closing_slide(prs, slide_index)

    return save_presentation(prs)


if __name__ == "__main__":
    output = build_presentation()
    print(output)
