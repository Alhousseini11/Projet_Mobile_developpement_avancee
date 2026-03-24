from __future__ import annotations

import math
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "presentation_assets" / "sprint_review"
OUTPUT_PATH = ROOT / "presentation_sprint_review_stories.pptx"

PRIMARY = RGBColor(220, 38, 38)
DARK = RGBColor(23, 31, 44)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(238, 242, 247)
WHITE = RGBColor(255, 255, 255)
SUCCESS = RGBColor(22, 101, 52)
SUCCESS_BG = RGBColor(220, 252, 231)
REVIEW = RGBColor(180, 83, 9)
REVIEW_BG = RGBColor(255, 237, 213)

PIL_PRIMARY = (220, 38, 38)
PIL_DARK = (15, 23, 42)
PIL_PANEL = (30, 41, 59)
PIL_LIGHT = (243, 246, 251)
PIL_TEXT = (226, 232, 240)
PIL_MUTED = (148, 163, 184)
PIL_WHITE = (255, 255, 255)

WIDE_W = 13.333
WIDE_H = 7.5

SCREENSHOTS = {
    "login": ROOT / "frontend" / "login-screen.png",
    "home": ROOT / "frontend" / "home-redesign.png",
    "profile": ROOT / "frontend" / "profile-redesign.png",
    "payment": ROOT / "frontend" / "payment-after-backend.png",
    "payment_alt": ROOT / "frontend" / "payment.png",
    "info": ROOT / "frontend" / "my-information.png",
    "invoices": ROOT / "frontend" / "invoices-screen.png",
    "post_login": ROOT / "frontend" / "post-login-screen.png",
    "vehicles": ROOT / "frontend" / "vehicles-final.png",
    "tutorials": ROOT / "frontend" / "tutorials-final.png",
}

LIVE_SCREENSHOT_FILES = {
    "reviews_live": ASSET_DIR / "reviews_live.png",
    "support_live": ASSET_DIR / "support_live.png",
}


STORIES = [
    {
        "number": "01",
        "title": "Reservation multi-service",
        "status": "Livree",
        "summary": [
            "Le mode multi permet de selectionner plusieurs prestations depuis un seul parcours de reservation.",
            "Le frontend boucle sur les services choisis et cree une reservation par service avec le meme vehicule, la meme date et le meme creneau.",
            "Le backend persiste chaque reservation creee, ce qui alimente ensuite le profil, les avis et les factures."
        ],
        "screens": ["reservation_multi"],
        "snippets": [
            {
                "path": "frontend/app/components/Reservations.vue",
                "anchor": "const serviceIdsToBook = multiServiceMode.value",
                "before": 0,
                "after": 28,
                "max_lines": 18,
                "label": "frontend/app/components/Reservations.vue"
            },
            {
                "path": "backend/src/modules/reservations/reservations.controller.ts",
                "anchor": "export const createReservation = async (req: Request, res: Response) => {",
                "before": 0,
                "after": 44,
                "max_lines": 18,
                "label": "backend/src/modules/reservations/reservations.controller.ts"
            }
        ]
    },
    {
        "number": "02",
        "title": "Avis et paiement",
        "status": "Livree",
        "summary": [
            "Paiement Stripe configure depuis le profil avec ouverture de Checkout puis synchro.",
            "Avis relies aux rendez-vous pour publier une note et un commentaire par intervention.",
            "Le backend persiste le moyen de paiement et les reviews par utilisateur."
        ],
        "screens": ["payment", "reviews_live"],
        "snippets": [
            {
                "path": "frontend/app/components/PaymentMethods.vue",
                "anchor": "async function startStripeSetup()",
                "before": 0,
                "after": 30,
                "max_lines": 18,
                "label": "frontend/app/components/PaymentMethods.vue"
            },
            {
                "path": "backend/src/modules/reviews/reviews.controller.ts",
                "anchor": "export async function upsertReview(req: Request, res: Response)",
                "before": 0,
                "after": 38,
                "max_lines": 18,
                "label": "backend/src/modules/reviews/reviews.controller.ts"
            }
        ]
    },
    {
        "number": "03",
        "title": "Support FAQ",
        "status": "Livree",
        "summary": [
            "Page support accessible depuis le profil avec contact atelier, horaires et FAQ.",
            "Questions frequentes ouvrables/fermables pour guider le client sans quitter l'app.",
            "Bloc message pret pour la demo et extensible vers un vrai support backend."
        ],
        "screens": ["support_live"],
        "snippets": [
            {
                "path": "frontend/app/components/SupportFAQ.vue",
                "anchor": "function toggleFaq(id: string)",
                "before": 4,
                "after": 26,
                "max_lines": 18,
                "label": "frontend/app/components/SupportFAQ.vue"
            }
        ]
    },
    {
        "number": "04",
        "title": "Bouton retour",
        "status": "Livree",
        "summary": [
            "Pattern de retour identique sur les ecrans profil, paiement, factures, FAQ et informations.",
            "La pile de navigation est geree de maniere centralisee pour eviter les impasses UX.",
            "Pendant la demo, l'utilisateur revient toujours a l'ecran precedent proprement."
        ],
        "screens": ["back_button_live"],
        "snippets": [
            {
                "path": "frontend/app/utils/navigation.ts",
                "anchor": "export function navigateToPage(page: AppPage, options: NavigationOptions = {})",
                "before": 0,
                "after": 32,
                "max_lines": 18,
                "label": "frontend/app/utils/navigation.ts"
            }
        ]
    },
    {
        "number": "05",
        "title": "Modification de rendez-vous",
        "status": "Livree",
        "summary": [
            "Depuis Mes rendez-vous, l'utilisateur peut reouvrir un RDV existant en edition.",
            "Le formulaire reservation recharge service, vehicule, date et creneau.",
            "Le backend reverifie les disponibilites avant de persister la mise a jour."
        ],
        "screens": ["post_login", "profile"],
        "snippets": [
            {
                "path": "frontend/app/components/Reservations.vue",
                "anchor": "async function updateReservation()",
                "before": 0,
                "after": 34,
                "max_lines": 18,
                "label": "frontend/app/components/Reservations.vue"
            },
            {
                "path": "backend/src/modules/reservations/reservations.controller.ts",
                "anchor": "export const updateReservation = async (req: Request, res: Response) => {",
                "before": 0,
                "after": 48,
                "max_lines": 18,
                "label": "backend/src/modules/reservations/reservations.controller.ts"
            }
        ]
    },
    {
        "number": "06",
        "title": "Modification information",
        "status": "Livree",
        "summary": [
            "Le client modifie ses coordonnees et preferences depuis Mes informations.",
            "Le payload est normalise cote frontend puis envoye a PUT /api/profile.",
            "La session reste coherente meme apres changement d'email."
        ],
        "screens": ["info"],
        "snippets": [
            {
                "path": "frontend/app/components/MyInformation.vue",
                "anchor": "function buildProfilePayload(): Partial<UserProfile> {",
                "before": 0,
                "after": 26,
                "max_lines": 18,
                "label": "frontend/app/components/MyInformation.vue"
            },
            {
                "path": "backend/src/modules/profile/profile.controller.ts",
                "anchor": "export async function updateProfile(req: Request, res: Response) {",
                "before": 0,
                "after": 48,
                "max_lines": 18,
                "label": "backend/src/modules/profile/profile.controller.ts"
            }
        ]
    },
    {
        "number": "07",
        "title": "Compteur de rendez-vous",
        "status": "Livree",
        "summary": [
            "Le profil affiche maintenant le nombre reel de reservations de l'utilisateur courant.",
            "Le calcul se fait cote backend a partir des reservations persistantes.",
            "Le frontend sait aussi resynchroniser le compteur pour eviter les zero incoherents."
        ],
        "screens": ["profile"],
        "snippets": [
            {
                "path": "frontend/app/services/ProfileService.ts",
                "anchor": "async function syncAppointmentCount(profile: UserProfile): Promise<UserProfile> {",
                "before": 0,
                "after": 22,
                "max_lines": 18,
                "label": "frontend/app/services/ProfileService.ts"
            },
            {
                "path": "backend/src/modules/profile/profile.controller.ts",
                "anchor": "async function buildProfileResponseForUser(",
                "before": 0,
                "after": 34,
                "max_lines": 18,
                "label": "backend/src/modules/profile/profile.controller.ts"
            }
        ]
    },
    {
        "number": "08",
        "title": "Bouton deconnexion",
        "status": "Livree",
        "summary": [
            "La deconnexion demande une confirmation avant fermeture de session.",
            "Le stockage local et l'etat reactif sont nettoyes en un seul point.",
            "L'app revient sur login avec historique vide pour proteger le compte."
        ],
        "screens": ["profile"],
        "snippets": [
            {
                "path": "frontend/app/components/Profile.vue",
                "anchor": "async function handleLogout() {",
                "before": 0,
                "after": 20,
                "max_lines": 18,
                "label": "frontend/app/components/Profile.vue"
            },
            {
                "path": "frontend/app/services/AuthService.ts",
                "anchor": "logout() {",
                "before": 0,
                "after": 10,
                "max_lines": 18,
                "label": "frontend/app/services/AuthService.ts"
            }
        ]
    },
    {
        "number": "09",
        "title": "Facture",
        "status": "Livree",
        "summary": [
            "Les factures sont listees par utilisateur et exposees au format PDF.",
            "Le mobile ouvre le PDF directement depuis la page Factures.",
            "Le backend derive les montants depuis les reservations et l'etat de paiement."
        ],
        "screens": ["invoices"],
        "snippets": [
            {
                "path": "frontend/app/components/Invoices.vue",
                "anchor": "async function openInvoicePdf(invoiceId: string) {",
                "before": 0,
                "after": 16,
                "max_lines": 18,
                "label": "frontend/app/components/Invoices.vue"
            },
            {
                "path": "backend/src/modules/profile/profile.controller.ts",
                "anchor": "return reservations.map(reservation => {",
                "before": 10,
                "after": 30,
                "max_lines": 18,
                "label": "backend/src/modules/profile/profile.controller.ts"
            }
        ]
    },
    {
        "number": "10",
        "title": "Notification",
        "status": "En review",
        "summary": [
            "Story encore en review sur ce sprint: le contrat API est cadre mais pas finalise.",
            "La route notifications est branchee dans l'API et renvoie pour l'instant un placeholder.",
            "Le prochain increment devra ajouter la vraie logique metier et l'UX mobile associee."
        ],
        "screens": ["home"],
        "snippets": [
            {
                "path": "backend/src/modules/notifications/notifications.controller.ts",
                "anchor": "export const listNotifications = createPlaceholderHandler('notifications', 'list');",
                "before": 1,
                "after": 1,
                "max_lines": 8,
                "label": "backend/src/modules/notifications/notifications.controller.ts"
            },
            {
                "path": "backend/src/modules/_shared/createPlaceholderHandler.ts",
                "anchor": "export function createPlaceholderHandler(feature: string, action: string) {",
                "before": 0,
                "after": 10,
                "max_lines": 18,
                "label": "backend/src/modules/_shared/createPlaceholderHandler.ts"
            }
        ]
    }
]


def ensure_asset_dir() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def load_font(size: int, mono: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    font_candidates = []
    if mono:
        font_candidates.extend(
            [
                Path("C:/Windows/Fonts/consola.ttf"),
                Path("C:/Windows/Fonts/consolab.ttf"),
                Path("C:/Windows/Fonts/cour.ttf"),
            ]
        )
    else:
        font_candidates.extend(
            [
                Path("C:/Windows/Fonts/segoeui.ttf"),
                Path("C:/Windows/Fonts/arial.ttf"),
                Path("C:/Windows/Fonts/calibri.ttf"),
            ]
        )

    for candidate in font_candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)

    return ImageFont.load_default()


def read_file_lines(relative_path: str) -> list[str]:
    return (ROOT / relative_path).read_text(encoding="utf-8").splitlines()


def trim_common_indent(lines: list[str]) -> list[str]:
    non_empty = [line for line in lines if line.strip()]
    if not non_empty:
        return lines

    indent = min(len(line) - len(line.lstrip(" ")) for line in non_empty)
    return [line[indent:] if len(line) >= indent else line for line in lines]


def extract_snippet(relative_path: str, anchor: str, before: int, after: int, max_lines: int) -> str:
    lines = read_file_lines(relative_path)
    anchor_index = next((index for index, line in enumerate(lines) if anchor in line), None)
    if anchor_index is None:
        raise ValueError(f"Anchor not found in {relative_path}: {anchor}")

    start = max(anchor_index - before, 0)
    end = min(anchor_index + after, len(lines) - 1)
    selected = trim_common_indent(lines[start : end + 1])[:max_lines]
    numbered = [f"{start + offset + 1:>4} | {line.rstrip()}" for offset, line in enumerate(selected)]
    return "\n".join(numbered)


def wrap_code_lines(snippet: str, max_chars: int) -> str:
    wrapped_lines: list[str] = []
    for raw_line in snippet.splitlines():
        if len(raw_line) <= max_chars:
            wrapped_lines.append(raw_line)
            continue

        prefix = raw_line[:7] if "|" in raw_line[:10] else ""
        content = raw_line[7:] if prefix else raw_line
        chunks = textwrap.wrap(
            content,
            width=max_chars - len(prefix),
            replace_whitespace=False,
            drop_whitespace=False,
            break_long_words=False,
            break_on_hyphens=False,
        )
        if not chunks:
            wrapped_lines.append(raw_line)
            continue

        wrapped_lines.append(f"{prefix}{chunks[0]}")
        continuation_prefix = " " * len(prefix)
        for chunk in chunks[1:]:
            wrapped_lines.append(f"{continuation_prefix}{chunk}")
    return "\n".join(wrapped_lines)


def render_code_panel(label: str, snippet: str, destination: Path) -> Path:
    width, height = 1700, 860
    image = Image.new("RGB", (width, height), PIL_DARK)
    draw = ImageDraw.Draw(image)
    title_font = load_font(32, mono=False)
    code_font = load_font(24, mono=True)
    small_font = load_font(22, mono=False)

    draw.rounded_rectangle((0, 0, width - 1, height - 1), radius=34, fill=PIL_DARK)
    draw.rounded_rectangle((0, 0, width - 1, 86), radius=34, fill=PIL_PANEL)
    draw.text((38, 22), label, font=title_font, fill=PIL_WHITE)
    draw.text((width - 235, 28), "code", font=small_font, fill=PIL_MUTED)

    max_chars = 92
    wrapped = wrap_code_lines(snippet, max_chars)
    top = 126
    line_height = 36
    for index, line in enumerate(wrapped.splitlines()):
        draw.text((40, top + index * line_height), line, font=code_font, fill=PIL_TEXT)

    image.save(destination)
    return destination


def render_reservation_multiservice_mock() -> Path:
    destination = ASSET_DIR / "reservation_multi_mock.png"
    width, height = 860, 1240
    image = Image.new("RGB", (width, height), (236, 240, 245))
    draw = ImageDraw.Draw(image)

    title_font = load_font(38)
    subtitle_font = load_font(24)
    label_font = load_font(22)
    value_font = load_font(24)
    strong_font = load_font(26)

    draw.rounded_rectangle((40, 34, width - 40, height - 34), radius=44, fill=PIL_WHITE)
    draw.rounded_rectangle((40, 34, width - 40, 138), radius=44, fill=PIL_DARK)
    draw.text((78, 66), "Reservation", font=title_font, fill=PIL_WHITE)
    draw.text((78, 106), "Mode multi-service active", font=subtitle_font, fill=(248, 113, 113))

    card_bounds = [
        (84, 188, width - 84, 348),
        (84, 374, width - 84, 534),
        (84, 560, width - 84, 720),
    ]
    cards = [
        ("Vidange", "45 min   79 CAD", True),
        ("Freins", "60 min   149 CAD", True),
        ("Diagnostic", "30 min   59 CAD", False),
    ]

    for bounds, (title, meta, selected) in zip(card_bounds, cards):
        x0, y0, x1, y1 = bounds
        fill = (255, 248, 248) if selected else PIL_WHITE
        outline = PIL_PRIMARY if selected else (203, 213, 225)
        draw.rounded_rectangle(bounds, radius=28, fill=fill, outline=outline, width=4)
        draw.text((x0 + 26, y0 + 22), title, font=strong_font, fill=PIL_DARK)
        draw.text((x0 + 26, y0 + 68), meta, font=label_font, fill=PIL_MUTED)
        draw.text((x0 + 26, y0 + 104), "Selectionne pour la meme date et le meme vehicule" if selected else "Disponible", font=label_font, fill=(71, 85, 105))
        badge_fill = PIL_PRIMARY if selected else (226, 232, 240)
        badge_text = PIL_WHITE if selected else PIL_PANEL
        draw.rounded_rectangle((x1 - 96, y0 + 52, x1 - 30, y0 + 112), radius=20, fill=badge_fill)
        draw.text((x1 - 74, y0 + 64), "OK" if selected else "+", font=strong_font, fill=badge_text)

    draw.rounded_rectangle((84, 760, width - 84, 930), radius=28, fill=(248, 250, 252))
    draw.text((112, 790), "Recapitulatif", font=strong_font, fill=PIL_DARK)
    draw.text((112, 838), "Services : Vidange • Freins", font=value_font, fill=PIL_PANEL)
    draw.text((112, 878), "Date : 24 Mars   Heure : 10:00", font=value_font, fill=PIL_PANEL)

    draw.rounded_rectangle((84, 968, width - 84, 1096), radius=28, fill=PIL_PRIMARY)
    draw.text((236, 1008), "Confirmer 2 services", font=strong_font, fill=PIL_WHITE)

    draw.text((118, 1140), "Mock visuel genere pour la slide story multi-service", font=subtitle_font, fill=PIL_MUTED)

    image.save(destination)
    return destination


def render_reviews_mock() -> Path:
    destination = ASSET_DIR / "reviews_mock.png"
    width, height = 860, 1240
    image = Image.new("RGB", (width, height), (238, 241, 245))
    draw = ImageDraw.Draw(image)

    title_font = load_font(34)
    strong_font = load_font(28)
    body_font = load_font(22)
    small_font = load_font(20)

    draw.rounded_rectangle((38, 28, width - 38, height - 28), radius=42, fill=PIL_WHITE)
    draw.rounded_rectangle((38, 28, width - 38, 128), radius=42, fill=PIL_DARK)
    draw.text((76, 62), "<", font=title_font, fill=PIL_WHITE)
    draw.text((124, 60), "Avis et evaluations", font=title_font, fill=PIL_WHITE)

    draw.rounded_rectangle((82, 156, width - 82, 266), radius=28, fill=PIL_WHITE, outline=(226, 232, 240), width=2)
    draw.text((112, 188), "Partagez votre retour atelier", font=strong_font, fill=PIL_DARK)
    draw.text((112, 226), "Notes, commentaires et suivi satisfaction client", font=body_font, fill=PIL_MUTED)

    draw.rounded_rectangle((82, 292, 394, 412), radius=26, fill=PIL_WHITE)
    draw.rounded_rectangle((466, 292, width - 82, 412), radius=26, fill=PIL_PANEL)
    draw.text((116, 326), "4.8", font=load_font(42), fill=PIL_DARK)
    draw.text((116, 370), "Note moyenne", font=body_font, fill=PIL_MUTED)
    draw.text((510, 326), "02", font=load_font(42), fill=PIL_WHITE)
    draw.text((510, 370), "Avis publies", font=body_font, fill=PIL_TEXT)

    draw.rounded_rectangle((82, 446, width - 82, 730), radius=30, fill=PIL_WHITE)
    draw.text((112, 478), "Deposer un avis", font=strong_font, fill=PIL_DARK)
    draw.rounded_rectangle((112, 528, width - 112, 612), radius=22, fill=(255, 248, 248), outline=PIL_PRIMARY, width=3)
    draw.text((138, 550), "Vidange", font=strong_font, fill=PIL_DARK)
    draw.text((138, 586), "24 mars 2026 a 10:00", font=small_font, fill=PIL_MUTED)
    draw.rounded_rectangle((width - 236, 548, width - 120, 600), radius=18, fill=(220, 252, 231))
    draw.text((width - 218, 562), "A evaluer", font=small_font, fill=(22, 101, 52))
    draw.text((112, 638), "★ ★ ★ ★ ☆", font=load_font(34), fill=(245, 158, 11))
    draw.rounded_rectangle((112, 680, width - 112, 720), radius=16, fill=(248, 250, 252))
    draw.text((132, 691), "Intervention claire, delai respecte...", font=small_font, fill=PIL_MUTED)

    draw.rounded_rectangle((82, 770, width - 82, 1120), radius=30, fill=PIL_WHITE)
    draw.text((112, 802), "Historique des avis", font=strong_font, fill=PIL_DARK)
    card_top = 848
    for title, stars, copy in [
        ("Vidange", "★★★★★", "Service rapide, accueil rassurant."),
        ("Freins", "★★★★☆", "Explications claires et atelier ponctuel."),
    ]:
        draw.rounded_rectangle((112, card_top, width - 112, card_top + 114), radius=20, fill=(248, 250, 252))
        draw.text((136, card_top + 18), title, font=body_font, fill=PIL_DARK)
        draw.text((width - 246, card_top + 18), stars, font=body_font, fill=(245, 158, 11))
        draw.text((136, card_top + 58), copy, font=small_font, fill=SLATE_TO_PIL())
        card_top += 132

    image.save(destination)
    return destination


def render_support_faq_mock() -> Path:
    destination = ASSET_DIR / "support_faq_mock.png"
    width, height = 860, 1240
    image = Image.new("RGB", (width, height), (236, 240, 245))
    draw = ImageDraw.Draw(image)

    title_font = load_font(34)
    strong_font = load_font(28)
    body_font = load_font(22)
    small_font = load_font(20)

    draw.rounded_rectangle((38, 28, width - 38, height - 28), radius=42, fill=PIL_WHITE)
    draw.rounded_rectangle((38, 28, width - 38, 128), radius=42, fill=PIL_DARK)
    draw.text((76, 62), "<", font=title_font, fill=PIL_WHITE)
    draw.text((124, 60), "Support et FAQ", font=title_font, fill=PIL_WHITE)

    draw.rounded_rectangle((84, 156, 258, 212), radius=24, fill=PIL_WHITE, outline=(226, 232, 240), width=2)
    draw.text((116, 172), "< Retour", font=body_font, fill=PIL_PANEL)

    draw.rounded_rectangle((84, 236, width - 84, 352), radius=28, fill=PIL_WHITE)
    draw.text((114, 270), "Aide et assistance", font=strong_font, fill=PIL_DARK)
    draw.text((114, 312), "FAQ, contact atelier et message client", font=body_font, fill=PIL_MUTED)

    draw.text((84, 392), "Contact atelier", font=strong_font, fill=PIL_DARK)
    y = 432
    for badge, title, desc, color in [
        ("TEL", "Appeler l'atelier", "+1 (438) 340 7583", (220, 38, 38)),
        ("@", "Envoyer un email", "contact@garageplus.ca", (14, 116, 144)),
        ("HRS", "Horaires d'ouverture", "Lun-Ven 8h-18h  |  Sam 9h-12h", (202, 138, 4)),
    ]:
        draw.rounded_rectangle((84, y, width - 84, y + 98), radius=24, fill=PIL_WHITE)
        draw.rounded_rectangle((108, y + 22, 178, y + 76), radius=16, fill=color)
        draw.text((123, y + 36), badge, font=small_font, fill=PIL_WHITE)
        draw.text((206, y + 26), title, font=body_font, fill=PIL_DARK)
        draw.text((206, y + 58), desc, font=small_font, fill=PIL_MUTED)
        y += 116

    draw.text((84, 812), "Questions frequentes", font=strong_font, fill=PIL_DARK)
    faq_blocks = [
        ("Comment prendre un rendez-vous ?", "Choisir un service, une date et confirmer."),
        ("Ou trouver mes factures ?", "Depuis Profil > Factures PDF."),
    ]
    y = 852
    for question, answer in faq_blocks:
        draw.rounded_rectangle((84, y, width - 84, y + 118), radius=24, fill=PIL_WHITE)
        draw.text((112, y + 20), question, font=body_font, fill=PIL_DARK)
        draw.text((112, y + 58), answer, font=small_font, fill=PIL_MUTED)
        draw.text((width - 132, y + 28), "▼", font=body_font, fill=PIL_MUTED)
        y += 136

    image.save(destination)
    return destination


def render_back_button_mock() -> Path:
    destination = ASSET_DIR / "back_button_mock.png"
    width, height = 860, 1240
    image = Image.new("RGB", (width, height), (236, 240, 245))
    draw = ImageDraw.Draw(image)

    title_font = load_font(34)
    strong_font = load_font(30)
    body_font = load_font(24)
    small_font = load_font(21)

    draw.rounded_rectangle((38, 28, width - 38, height - 28), radius=42, fill=PIL_WHITE)
    draw.rounded_rectangle((38, 28, width - 38, 128), radius=42, fill=PIL_DARK)
    draw.text((76, 62), "<", font=title_font, fill=PIL_WHITE)
    draw.text((124, 60), "Demonstration retour", font=title_font, fill=PIL_WHITE)

    draw.rounded_rectangle((84, 170, 292, 238), radius=28, fill=PIL_WHITE, outline=PIL_PRIMARY, width=4)
    draw.text((116, 190), "< Retour", font=strong_font, fill=PIL_PRIMARY)
    draw.text((84, 264), "Le bouton retour est volontairement mis en avant ici.", font=small_font, fill=PIL_MUTED)

    draw.rounded_rectangle((84, 330, width - 84, 520), radius=30, fill=(255, 248, 248))
    draw.text((116, 366), "Pattern commun", font=strong_font, fill=PIL_DARK)
    draw.text((116, 414), "ActionBar: fleche en haut a gauche", font=body_font, fill=PIL_PANEL)
    draw.text((116, 452), "Inline back: pillule blanche sous l'entete", font=body_font, fill=PIL_PANEL)

    draw.rounded_rectangle((84, 574, width - 84, 1050), radius=30, fill=PIL_WHITE)
    draw.text((116, 612), "Ecrans concernes", font=strong_font, fill=PIL_DARK)
    for index, line in enumerate([
        "Profil",
        "Mes informations",
        "Factures PDF",
        "Mode de paiement",
        "Avis et evaluations",
        "Support et FAQ",
        "Reservations / Modification de rendez-vous",
    ]):
        top = 668 + index * 52
        draw.rounded_rectangle((120, top, 144, top + 24), radius=12, fill=PIL_PRIMARY)
        draw.text((160, top - 2), line, font=body_font, fill=PIL_PANEL)

    draw.rounded_rectangle((84, 1088, width - 84, 1148), radius=22, fill=(220, 252, 231))
    draw.text((118, 1104), "Le visuel montre explicitement le bouton retour, sans qu'il se perde dans la slide.", font=small_font, fill=(22, 101, 52))

    image.save(destination)
    return destination


def render_back_button_focus_live(source_path: Path) -> Path:
    destination = ASSET_DIR / "back_button_live.png"
    width, height = 860, 1240
    canvas = Image.new("RGB", (width, height), (236, 240, 245))
    draw = ImageDraw.Draw(canvas)

    with Image.open(source_path) as source:
        screenshot = source.convert("RGB")
        full_preview = ImageOps.fit(screenshot, (360, 640), method=Image.Resampling.LANCZOS)
        crop = screenshot.crop((0, 32, 720, 250))
        crop_preview = ImageOps.fit(crop, (740, 360), method=Image.Resampling.LANCZOS)

    draw.rounded_rectangle((42, 32, width - 42, 458), radius=36, fill=PIL_WHITE)
    canvas.paste(crop_preview, (60, 50))
    draw.rounded_rectangle((78, 120, 330, 204), radius=22, outline=PIL_PRIMARY, width=6)
    draw.rounded_rectangle((42, 496, width - 42, height - 42), radius=36, fill=PIL_WHITE)
    canvas.paste(full_preview, (250, 540))

    title_font = load_font(30)
    body_font = load_font(22)
    draw.text((70, 426), "Retour visible: ActionBar + bouton inline", font=title_font, fill=PIL_DARK)
    draw.text((70, 1160), "Capture reelle de l'application sur emulateur", font=body_font, fill=PIL_MUTED)

    canvas.save(destination)
    return destination


def SLATE_TO_PIL() -> tuple[int, int, int]:
    return (71, 85, 105)


def fit_image(path: Path, size: tuple[int, int]) -> Image.Image:
    with Image.open(path) as source:
        return ImageOps.fit(source.convert("RGB"), size, method=Image.Resampling.LANCZOS)


def draw_shadow(draw: ImageDraw.ImageDraw, bounds: tuple[int, int, int, int]) -> None:
    x0, y0, x1, y1 = bounds
    shadow = (x0 + 8, y0 + 12, x1 + 8, y1 + 12)
    draw.rounded_rectangle(shadow, radius=34, fill=(207, 216, 226))


def paste_card(base: Image.Image, screenshot: Image.Image, bounds: tuple[int, int, int, int]) -> None:
    draw = ImageDraw.Draw(base)
    draw_shadow(draw, bounds)
    x0, y0, x1, y1 = bounds
    draw.rounded_rectangle(bounds, radius=34, fill=PIL_WHITE)
    card = Image.new("RGB", (x1 - x0 - 36, y1 - y0 - 36), PIL_WHITE)
    card.paste(screenshot, (0, 0))
    mask = Image.new("L", card.size, 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, card.size[0], card.size[1]), radius=26, fill=255)
    base.paste(card, (x0 + 18, y0 + 18), mask)


def render_visual_collage(name: str, screen_keys: list[str]) -> Path:
    destination = ASSET_DIR / f"{name}_visual.png"
    width, height = 1100, 1560
    image = Image.new("RGB", (width, height), PIL_LIGHT)
    draw = ImageDraw.Draw(image)

    photos = [fit_image(SCREENSHOTS[key], (860, 1240)) for key in screen_keys]

    if len(photos) == 1:
        paste_card(image, photos[0], (90, 100, width - 90, height - 90))
    elif len(photos) == 2:
        slots = [
            (86, 80, width - 86, 770),
            (86, 790, width - 86, height - 80),
        ]
        for photo, bounds in zip(photos, slots):
            sized = ImageOps.fit(photo, (bounds[2] - bounds[0] - 36, bounds[3] - bounds[1] - 36), method=Image.Resampling.LANCZOS)
            paste_card(image, sized, bounds)
    else:
        slots = [
            (60, 80, 525, 760),
            (575, 80, 1040, 760),
            (60, 810, 525, 1490),
            (575, 810, 1040, 1490),
        ]
        for photo, bounds in zip(photos, slots):
            sized = ImageOps.fit(photo, (bounds[2] - bounds[0] - 36, bounds[3] - bounds[1] - 36), method=Image.Resampling.LANCZOS)
            paste_card(image, sized, bounds)

    image.save(destination)
    return destination


def build_story_assets() -> dict[str, dict[str, object]]:
    assets: dict[str, dict[str, object]] = {}
    for story in STORIES:
        visual = render_visual_collage(f"story_{story['number']}", story["screens"])
        code_panels: list[Path] = []
        for panel_index, snippet_config in enumerate(story["snippets"], start=1):
            snippet = extract_snippet(
                snippet_config["path"],
                snippet_config["anchor"],
                snippet_config["before"],
                snippet_config["after"],
                snippet_config["max_lines"],
            )
            code_path = ASSET_DIR / f"story_{story['number']}_code_{panel_index}.png"
            code_panels.append(render_code_panel(snippet_config["label"], snippet, code_path))
        assets[story["number"]] = {
            "visual": visual,
            "code_panels": code_panels,
        }
    return assets


def add_full_slide_background(slide) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str,
                font_size: int, color: RGBColor = DARK, bold: bool = False,
                font_name: str = "Segoe UI", align=PP_ALIGN.LEFT) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    paragraph.alignment = align


def add_status_chip(slide, status: str) -> None:
    left, top, width, height = 11.25, 0.46, 1.45, 0.44
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()

    if status.lower().startswith("en review"):
        fill.fore_color.rgb = REVIEW_BG
        text_color = REVIEW
    else:
        fill.fore_color.rgb = SUCCESS_BG
        text_color = SUCCESS

    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = status
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Segoe UI"
    paragraph.alignment = PP_ALIGN.CENTER


def add_bullet_box(slide, left: float, top: float, width: float, height: float, items: list[str]) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.line.fill.background()
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.2)
    frame.margin_right = Inches(0.15)
    frame.margin_top = Inches(0.15)
    frame.margin_bottom = Inches(0.1)

    title = frame.paragraphs[0]
    title_run = title.add_run()
    title_run.text = "Points cles"
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = DARK
    title_run.font.name = "Segoe UI"

    for item in items:
        paragraph = frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.level = 0
        paragraph.font.size = Pt(15)
        paragraph.font.color.rgb = SLATE
        paragraph.font.name = "Segoe UI"
        paragraph.space_before = Pt(4)


def add_story_slide(prs: Presentation, story: dict[str, object], assets: dict[str, dict[str, object]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_background(slide)

    story_number = story["number"]
    visual_path = assets[story_number]["visual"]
    code_panels: list[Path] = assets[story_number]["code_panels"]

    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.56),
        Inches(0.46),
        Inches(1.02),
        Inches(0.44),
    )
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = PRIMARY
    chip_frame = chip.text_frame
    chip_frame.clear()
    chip_p = chip_frame.paragraphs[0]
    chip_run = chip_p.add_run()
    chip_run.text = f"Story {story_number}"
    chip_run.font.size = Pt(15)
    chip_run.font.bold = True
    chip_run.font.color.rgb = WHITE
    chip_run.font.name = "Segoe UI"
    chip_p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 1.72, 0.34, 8.9, 0.52, str(story["title"]), 26, bold=True)
    add_status_chip(slide, str(story["status"]))
    add_textbox(
        slide,
        0.58,
        1.02,
        11.9,
        0.28,
        "Sprint review - application mobile garage et API backend",
        11,
        color=SLATE,
    )

    slide.shapes.add_picture(str(visual_path), Inches(0.62), Inches(1.42), width=Inches(3.74), height=Inches(5.52))
    add_bullet_box(slide, 4.62, 1.42, 8.08, 1.45, list(story["summary"]))

    if len(code_panels) == 1:
        slide.shapes.add_picture(str(code_panels[0]), Inches(4.62), Inches(3.12), width=Inches(8.08), height=Inches(3.72))
    else:
        slide.shapes.add_picture(str(code_panels[0]), Inches(4.62), Inches(3.12), width=Inches(3.92), height=Inches(3.72))
        slide.shapes.add_picture(str(code_panels[1]), Inches(8.78), Inches(3.12), width=Inches(3.92), height=Inches(3.72))

    add_textbox(
        slide,
        0.62,
        7.02,
        12.0,
        0.22,
        "Presentation generee automatiquement a partir des composants, services et controllers du repo.",
        9,
        color=SLATE,
    )


def add_title_slide(prs: Presentation, cover_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_background(slide)

    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.64),
        Inches(0.56),
        Inches(1.66),
        Inches(0.46),
    )
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = PRIMARY
    chip_frame = chip.text_frame
    chip_frame.clear()
    chip_p = chip_frame.paragraphs[0]
    chip_run = chip_p.add_run()
    chip_run.text = "Sprint review"
    chip_run.font.size = Pt(16)
    chip_run.font.bold = True
    chip_run.font.color.rgb = WHITE
    chip_run.font.name = "Segoe UI"
    chip_p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 0.64, 1.2, 6.8, 0.9, "Presentation des stories du sprint", 30, bold=True)
    add_textbox(
        slide,
        0.64,
        2.02,
        5.8,
        1.0,
        "Stories livrees du sprint, reservation multi-service, parcours client et point en review pour la demo.",
        18,
        color=SLATE,
    )
    add_textbox(slide, 0.64, 2.82, 5.0, 0.28, "Mars 2026 | Projet Mobile Developpement Avancee", 11, color=SLATE)
    slide.shapes.add_picture(str(cover_path), Inches(6.0), Inches(0.84), width=Inches(6.65), height=Inches(5.95))

    footer = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.64),
        Inches(6.46),
        Inches(12.0),
        Inches(0.5),
    )
    footer.line.fill.background()
    footer.fill.solid()
    footer.fill.fore_color.rgb = WHITE
    footer_frame = footer.text_frame
    footer_frame.clear()
    footer_p = footer_frame.paragraphs[0]
    footer_run = footer_p.add_run()
    footer_run.text = "Une slide par story, avec ecran de l'application et extrait de code explicatif."
    footer_run.font.size = Pt(14)
    footer_run.font.color.rgb = DARK
    footer_run.font.name = "Segoe UI"
    footer_p.alignment = PP_ALIGN.CENTER


def build_cover_collage() -> Path:
    destination = ASSET_DIR / "cover_collage.png"
    return render_visual_collage("cover", ["login", "payment", "info", "invoices"])


def save_presentation(prs: Presentation) -> Path:
    candidates = [OUTPUT_PATH]
    candidates.extend(
        OUTPUT_PATH.with_name(f"{OUTPUT_PATH.stem}_v{index}{OUTPUT_PATH.suffix}")
        for index in range(2, 10)
    )

    last_error: PermissionError | None = None
    for candidate in candidates:
        try:
            prs.save(candidate)
            return candidate
        except PermissionError as error:
            last_error = error
            continue

    raise last_error or PermissionError(f"Unable to save presentation: {OUTPUT_PATH}")


def select_screenshot(key: str, fallback_key: str) -> Path:
    live_path = LIVE_SCREENSHOT_FILES.get(key)
    if live_path and live_path.exists():
        return live_path
    return SCREENSHOTS[fallback_key]


def build_presentation() -> Path:
    ensure_asset_dir()
    SCREENSHOTS["reservation_multi"] = render_reservation_multiservice_mock()
    SCREENSHOTS["reviews_mock"] = render_reviews_mock()
    SCREENSHOTS["support_faq_mock"] = render_support_faq_mock()
    SCREENSHOTS["back_button_mock"] = render_back_button_mock()
    SCREENSHOTS["reviews_live"] = select_screenshot("reviews_live", "reviews_mock")
    SCREENSHOTS["support_live"] = select_screenshot("support_live", "support_faq_mock")
    if LIVE_SCREENSHOT_FILES["support_live"].exists():
        SCREENSHOTS["back_button_live"] = render_back_button_focus_live(LIVE_SCREENSHOT_FILES["support_live"])
    else:
        SCREENSHOTS["back_button_live"] = SCREENSHOTS["back_button_mock"]
    story_assets = build_story_assets()
    cover_path = build_cover_collage()

    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)

    add_title_slide(prs, cover_path)
    for story in STORIES:
        add_story_slide(prs, story, story_assets)

    return save_presentation(prs)


if __name__ == "__main__":
    output = build_presentation()
    print(output)
